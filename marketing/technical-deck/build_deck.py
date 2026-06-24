#!/usr/bin/env python3
"""
Build a Google-Slides-compatible .pptx *technical* deck from the content of
.smithers/workflows/demo.tsx (the keyboard-driven terminal demo that
.smithers/scripts/run-demo.sh plays).

Output: smithers-technical-deck.pptx — import into Google Slides via
"File -> Import slides" or just upload to Drive and "Open with Google Slides".

Every slide's spoken narration from the terminal deck is preserved verbatim as
PowerPoint speaker notes, so the deck is presentation-ready. The terminal look
(dark panels, monospace, colored "boxes") is reproduced with native, editable
shapes and text — no flat screenshots.

Run:  marketing/technical-deck/.venv/bin/python marketing/technical-deck/build_deck.py
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "smithers-technical-deck.pptx")


# ── palette (mirrors the ANSI colors used in demo.tsx) ───────────────────────
def rgb(h):
    return RGBColor.from_string(h)


BG = rgb("0E131C")
PANEL = rgb("161D2A")
PANEL2 = rgb("1B2433")
CODEBG = rgb("0B1119")
BORDER = rgb("2A3445")
TEXT = rgb("E6EDF3")
DIM = rgb("93A1B3")
MUTED = rgb("64748B")
WHITE = rgb("FFFFFF")
CYAN = rgb("56C7E6")
TEAL = rgb("78D7BA")
GREEN = rgb("4ADE80")
YELLOW = rgb("FBBF24")
MAGENTA = rgb("C084FC")
RED = rgb("F87171")
REDSOFT = rgb("FCA5A5")

HEAD = "Inter"
BODY = "Inter"
MONO = "Roboto Mono"

prs = Presentation()
prs.slide_width = In(13.333)
prs.slide_height = In(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

# body region for stacked content
LX = 0.6
LW = 12.13
BODY_TOP = 2.45
BODY_BOT = 7.16
AVAIL = BODY_BOT - BODY_TOP


# ── primitives ───────────────────────────────────────────────────────────────
def seg(t, color=TEXT, size=12, bold=False, font=BODY, italic=False):
    return {"t": t, "color": color, "size": size, "bold": bold, "font": font, "italic": italic}


def _no_shadow(sh):
    sh.shadow.inherit = False


def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = slide.shapes.add_shape(shape, In(l), In(t), In(w), In(h))
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    _no_shadow(sh)
    return sh


def card(slide, l, t, w, h, fill=PANEL, line=BORDER, radius=0.04, line_w=1.0):
    return rect(slide, l, t, w, h, fill=fill, line=line, line_w=line_w,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)


def para(tf, segs, align=PP_ALIGN.LEFT, space_after=0, line_spacing=1.0, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    for s in segs:
        r = p.add_run()
        r.text = s["t"]
        f = r.font
        f.size = Pt(s["size"])
        f.bold = s["bold"]
        f.italic = s["italic"]
        f.name = s["font"]
        f.color.rgb = s["color"]
    return p


def txt(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=0, line_spacing=1.0, wrap=True):
    tb = slide.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, segs in enumerate(lines):
        para(tf, segs, align=align, space_after=space_after, line_spacing=line_spacing, first=(i == 0))
    return tb


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ── code syntax painter (mirrors printCode in demo.tsx) ──────────────────────
CODE_KW = re.compile(r"\b(import|from|export|default|const|let|return|function|async|await|"
                     r"if|else|for|of|new|class|interface|type)\b")


def paint_code_line(line):
    """Return a list of (text, color) runs for one line of code."""
    n = len(line)
    colors = [TEXT] * n
    locked = [False] * n

    def apply(rx, color, lock=False):
        for m in re.finditer(rx, line):
            for i in range(m.start(), m.end()):
                if not locked[i]:
                    colors[i] = color
                    if lock:
                        locked[i] = True

    apply(r"//.*", MUTED, lock=True)                                         # comments
    apply(r"\"(?:[^\"\\]|\\.)*\"|'(?:[^'\\]|\\.)*'|`(?:[^`\\]|\\.)*`", GREEN, lock=True)  # strings
    apply(CODE_KW.pattern, MAGENTA)                                          # keywords
    apply(r"</?\w+", CYAN)                                                   # JSX tags
    apply(r"\b\w+(?==)", YELLOW)                                             # prop names

    runs, cur, col = [], "", (colors[0] if n else TEXT)
    for ch, c in zip(line, colors):
        if c == col:
            cur += ch
        else:
            runs.append((cur, col))
            cur, col = ch, c
    if cur or not runs:
        runs.append((cur, col))
    return runs


def _trim(lines):
    while lines and lines[0].strip() == "":
        lines = lines[1:]
    while lines and lines[-1].strip() == "":
        lines = lines[:-1]
    return lines


def code_size(text, w, max_h):
    lines = _trim(text.split("\n"))
    n = max(1, len(lines))
    maxlen = max((len(l) for l in lines), default=1)
    f_h = 60.0 * (max_h - 0.46) / n          # line-spacing factor 1.2
    f_w = (w - 0.5) * 72.0 / (0.6 * maxlen)  # keep longest line on one line
    f = min(13.0, f_h, f_w)
    return max(8.0, f), lines


def code_panel(slide, l, t, w, text, file=None, size=None, max_h=None):
    if size is None:
        size, lines = code_size(text, w, max_h if max_h else AVAIL)
    else:
        lines = _trim(text.split("\n"))
    lh = 1.2 * size / 72.0
    head = 0.30
    h = head + len(lines) * lh + 0.16
    card(slide, l, t, w, h, fill=CODEBG, line=BORDER, radius=0.03)
    # window chrome
    for i, dot in enumerate((rgb("F87171"), rgb("FBBF24"), rgb("4ADE80"))):
        rect(slide, l + 0.22 + i * 0.22, t + 0.12, 0.085, 0.085, fill=dot, shape=MSO_SHAPE.OVAL)
    if file:
        txt(slide, l + 1.0, t + 0.04, w - 1.2, 0.24,
            [[seg(file, MUTED, 9.5, font=MONO)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, l, t + head - 0.02, w, 0.012, fill=rgb("141C28"))
    y = t + head
    for line in lines:
        runs = paint_code_line(line.replace("\t", "  "))
        txt(slide, l + 0.26, y, w - 0.4, lh + 0.04,
            [[seg(tx, co, size, font=MONO) for tx, co in runs]], wrap=False, line_spacing=1.0)
        y += lh
    return h


def code_panel_h(text, w, size=None, max_h=None):
    if size is None:
        size, lines = code_size(text, w, max_h if max_h else AVAIL)
    else:
        lines = _trim(text.split("\n"))
    lh = 1.2 * size / 72.0
    return 0.30 + len(lines) * lh + 0.16


# ── terminal-style box (mirrors box() in demo.tsx) ───────────────────────────
# A "line" is one of:
#   "string"                         -> dim body text, full width
#   [seg, seg, ...]                  -> a pre-built run list
#   ("kv", key, val, color)          -> aligned key/value row (key mono+bold)
#   ("kv", key, val, color, ksize)
def _is_blank(line):
    return isinstance(line, str) and line.strip() == ""


def box_line_h(size):
    return size * 0.026 + 0.02


def box_h(b):
    size = b.get("size", 12)
    lh = box_line_h(size)
    head = 0.5 if b.get("title") else 0.22
    body = sum((lh * (0.5 if _is_blank(l) else 1.0)) for l in b["lines"])
    return head + body + 0.18


def draw_box(slide, l, t, w, b, h=None):
    size = b.get("size", 12)
    color = b.get("color", CYAN)
    lh = box_line_h(size)
    if h is None:
        h = box_h(b)
    padx = 0.32
    card(slide, l, t, w, h, fill=PANEL, line=color, radius=0.04, line_w=1.25)
    rect(slide, l, t, w, 0.06, fill=color)
    y = t + 0.18
    if b.get("title"):
        txt(slide, l + padx, y, w - 2 * padx, 0.34,
            [[seg(b["title"], color, size + 1.5, bold=True, font=HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.30
        rect(slide, l + padx, y, w - 2 * padx, 0.011, fill=BORDER)
        y += 0.12
    else:
        y += 0.02
    kvw = b.get("kvw", 2.5)
    for line in b["lines"]:
        if _is_blank(line):
            y += lh * 0.5
            continue
        if isinstance(line, tuple) and line and line[0] == "kv":
            _, key, val, kcolor = line[0], line[1], line[2], line[3]
            ksize = line[4] if len(line) > 4 else size
            txt(slide, l + padx, y, kvw, lh + 0.05,
                [[seg(key, kcolor, ksize, bold=True, font=MONO)]], anchor=MSO_ANCHOR.MIDDLE)
            txt(slide, l + padx + kvw, y, w - 2 * padx - kvw, lh + 0.05,
                [[seg(val, DIM, size, font=BODY)]], anchor=MSO_ANCHOR.MIDDLE)
        elif isinstance(line, str):
            txt(slide, l + padx, y, w - 2 * padx, lh + 0.05,
                [[seg(line, DIM, size, font=BODY)]], anchor=MSO_ANCHOR.MIDDLE)
        else:  # seg list
            txt(slide, l + padx, y, w - 2 * padx, lh + 0.05, [line], anchor=MSO_ANCHOR.MIDDLE)
        y += lh
    return h


# ── slide chrome ─────────────────────────────────────────────────────────────
def base(tag, accent, title, subtitle, page, total):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=BG)
    rect(s, 0, 0, SW, 0.06, fill=rgb("0B0F17"))
    txt(s, 0.6, 0.30, 7, 0.3,
        [[seg("smithers", WHITE, 11, bold=True, font=MONO), seg("  ·  technical deck", DIM, 11, font=MONO)]])
    txt(s, SW - 3.1, 0.30, 2.5, 0.3, [[seg(f"{page:02d} / {total}", DIM, 11, font=MONO)]], align=PP_ALIGN.RIGHT)
    if tag:
        rect(s, 0.6, 0.86, 0.16, 0.16, fill=accent)
        txt(s, 0.86, 0.80, 10, 0.3, [[seg(tag, accent, 12, bold=True, font=MONO)]])
    txt(s, 0.58, 1.14, 12.1, 0.85, [[seg(title, WHITE, 29, bold=True, font=HEAD)]], line_spacing=1.0)
    if subtitle:
        txt(s, 0.6, 1.82, 12.1, 0.45, [[seg(subtitle, DIM, 14.5, font=BODY)]])
    rect(s, 0.6, 2.28 if subtitle else 1.78, 1.7, 0.035, fill=accent)
    return s


# ── special block renderers ──────────────────────────────────────────────────
def draw_runtime(slide, l, t, w):
    nodes = [("Render", "tree"), ("Extract", "tasks"), ("Execute", "ready"), ("Persist", "outputs")]
    nw, gap = 2.55, 0.5
    nh = 0.95
    x = l
    for i, (a, bl) in enumerate(nodes):
        card(slide, x, t, nw, nh, fill=PANEL, line=CYAN, radius=0.08, line_w=1.4)
        txt(slide, x, t + 0.14, nw, 0.4, [[seg(a, CYAN, 16, bold=True, font=MONO)]],
            align=PP_ALIGN.CENTER)
        txt(slide, x, t + 0.52, nw, 0.3, [[seg(bl, DIM, 11.5, font=MONO)]], align=PP_ALIGN.CENTER)
        if i < 3:
            txt(slide, x + nw, t, gap, nh, [[seg("→", MUTED, 20, bold=True, font=MONO)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += nw + gap
    span = nw * 4 + gap * 3
    # return loop
    rect(slide, l + 0.4, t + nh + 0.42, span - 0.8, 0.022, fill=MUTED)
    rect(slide, l + 0.4, t + nh + 0.06, 0.022, 0.38, fill=MUTED)
    rect(slide, l + span - 0.4, t + nh + 0.06, 0.022, 0.38, fill=MUTED)
    txt(slide, l, t + nh + 0.50, span, 0.3,
        [[seg("re-render with new state", DIM, 12, font=MONO)]], align=PP_ALIGN.CENTER)
    by = t + nh + 0.95
    bx = {"title": "Why one-way data flow matters", "color": MAGENTA, "size": 12, "lines": [
        [seg("Events update state. State is the source of truth. The plan is a ", TEXT, 12.5, font=BODY),
         seg("pure function", WHITE, 12.5, bold=True, font=BODY), seg(" of state.", TEXT, 12.5, font=BODY)],
        "",
        [seg("Free time travel  ", MAGENTA, 12, bold=True, font=MONO),
         seg("a frame is a snapshot; forking is “throw away rows”", DIM, 12, font=BODY)],
        [seg("Free resume       ", MAGENTA, 12, bold=True, font=MONO),
         seg("re-render from current state, no event log to replay", DIM, 12, font=BODY)],
        [seg("Free SQL debug    ", MAGENTA, 12, bold=True, font=MONO),
         seg("state is queryable, an event chain is not", DIM, 12, font=BODY)],
    ]}
    draw_box(slide, l, by, w, bx)
    return (by + box_h(bx)) - t


def runtime_h(w):
    return 3.95


def draw_layers(slide, l, t, w):
    rows = [
        (RED, "MODEL LAYER", "GPT · Claude · Gemini · Kimi — whatever wins this week", "volatile · weekly", False),
        (YELLOW, "AGENT / TOPOLOGY LAYER", "ReAct · crew · swarm · plan-execute · background", "fluid · quarterly", False),
        (GREEN, "ORCHESTRATION LAYER — this is Smithers", "durable steps · retries · state · events · observability", "stable · build once", True),
    ]
    lh = 1.12
    y = t
    for col, name, detail, speed, hot in rows:
        card(slide, l, y, w, lh, fill=(rgb("13271B") if hot else PANEL),
             line=(col if hot else BORDER), radius=0.05, line_w=(1.75 if hot else 1.0))
        rect(slide, l, y, 0.12, lh, fill=col)
        txt(slide, l + 0.4, y + 0.20, w - 4.0, 0.4,
            [[seg(name, col if hot else TEXT, 16.5 if hot else 15.5, bold=True, font=HEAD)]])
        txt(slide, l + 0.4, y + 0.64, w - 4.0, 0.4, [[seg(detail, DIM, 12.5, font=MONO)]])
        txt(slide, l + w - 3.4, y, 3.0, lh, [[seg(speed, col, 12.5, bold=hot, font=BODY)]],
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y += lh + 0.18
    return y - t - 0.18


def layers_h(w):
    return 1.12 * 3 + 0.18 * 2


# ── block dispatch ───────────────────────────────────────────────────────────
def block_h(b, w):
    ty = b["t"]
    if ty == "box":
        return box_h(b)
    if ty == "code":
        return code_panel_h(b["text"], w, size=b.get("size"), max_h=b.get("max_h"))
    if ty == "cols":
        cw = (w - 0.24) / 2
        return max(box_h(b["items"][0]), box_h(b["items"][1]))
    if ty == "code2":
        cw = (w - 0.24) / 2
        return max(code_panel_h(b["left"], cw, size=b.get("size")),
                   code_panel_h(b["right"], cw, size=b.get("size")))
    if ty == "split":
        return AVAIL - 0.1
    if ty == "runtime":
        return runtime_h(w)
    if ty == "layers":
        return layers_h(w)
    if ty == "para":
        n = len(b["lines"])
        return n * (b.get("size", 13) * 0.027 + 0.04) + 0.1
    if ty == "gap":
        return b.get("h", 0.2)
    return 0.0


def draw_block(slide, b, t, w):
    ty = b["t"]
    if ty == "box":
        return draw_box(slide, LX, t, w, b)
    if ty == "code":
        return code_panel(slide, LX, t, w, b["text"], file=b.get("file"),
                          size=b.get("size"), max_h=b.get("max_h"))
    if ty == "cols":
        cw = (w - 0.24) / 2
        h = max(box_h(b["items"][0]), box_h(b["items"][1]))
        draw_box(slide, LX, t, cw, b["items"][0], h=h)
        draw_box(slide, LX + cw + 0.24, t, cw, b["items"][1], h=h)
        return h
    if ty == "code2":
        cw = (w - 0.24) / 2
        size = b.get("size")
        h = max(code_panel_h(b["left"], cw, size=size), code_panel_h(b["right"], cw, size=size))
        code_panel(slide, LX, t, cw, b["left"], file=b.get("lfile"), size=size)
        code_panel(slide, LX + cw + 0.24, t, cw, b["right"], file=b.get("rfile"), size=size)
        return h
    if ty == "split":
        lw = b.get("lw", 7.15)
        rw = w - lw - 0.25
        h = AVAIL - 0.1
        code_panel(slide, LX, t, lw, b["code"], file=b.get("file"), max_h=h)
        draw_box(slide, LX + lw + 0.25, t, rw, b["box"], h=h)
        return h
    if ty == "runtime":
        return draw_runtime(slide, LX, t, w)
    if ty == "layers":
        return draw_layers(slide, LX, t, w)
    if ty == "para":
        h = block_h(b, w)
        txt(slide, LX, t, w, h, b["lines"], align=b.get("align", PP_ALIGN.LEFT),
            space_after=4, line_spacing=1.18)
        return h
    if ty == "gap":
        return b.get("h", 0.2)
    return 0.0


def render_content(s, blocks):
    gap = 0.22
    heights = [block_h(b, LW) for b in blocks]
    total = sum(heights) + gap * (len(blocks) - 1)
    y = BODY_TOP + max(0.0, (AVAIL - total) / 2.0)
    for b, h in zip(blocks, heights):
        drawn = draw_block(s, b, y, LW)
        y += (drawn if drawn else h) + gap


# helper shortcuts for authoring box lines
def kv(k, v, c, ksize=None):
    return ("kv", k, v, c) if ksize is None else ("kv", k, v, c, ksize)


def mono(t, c=CYAN, size=12, bold=True):
    return seg(t, c, size, bold=bold, font=MONO)


def dim(t, size=12):
    return seg(t, DIM, size, font=BODY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════════════════════════
SLIDES = []


def slide(**kw):
    SLIDES.append(kw)


# 01 — TITLE ------------------------------------------------------------------
slide(kind="title",
      notes="Hi. This is Smithers. Over the next few minutes I'm going to walk through the problem it "
            "solves, the full set of features it ships with, and a couple of live demos of it doing "
            "things a queue and a database can't.")

# 02 — THE TREADMILL ----------------------------------------------------------
slide(tag="ACT I · THE PROBLEM", accent=RED, title="The treadmill",
      subtitle="Every six months, the right way to build an AI agent changes.",
      notes="Every six months, the right way to build an AI agent changes. Chains. ReAct. Tools. Plan "
            "and execute. Multi-agent. Crews. Swarms. If you coupled your infrastructure to any one of "
            "these, you've already rebuilt twice. And you'll rebuild again.",
      blocks=[
          {"t": "box", "title": "Six months apart, every time", "color": YELLOW, "size": 13, "kvw": 1.8, "lines": [
              kv("2023 Q1", "chains", MUTED),
              kv("2023 Q3", "ReAct", MUTED),
              kv("2024 Q1", "tools / function calling", MUTED),
              kv("2024 Q3", "plan-and-execute, planner + worker", MUTED),
              kv("2025 Q1", "multi-agent / crews / swarms", MUTED),
              kv("2025 Q3", "background agents", MUTED),
              kv("2026 Q1", "? (the meta keeps moving)", MUTED),
              "",
              [seg("If your infra is coupled to the topology, you rebuild every time.", YELLOW, 13.5, bold=True, font=BODY)],
          ]},
      ])

# 03 — BACKGROUND AGENTS ARE DIFFERENT ----------------------------------------
slide(tag="ACT I · THE PROBLEM", accent=RED, title="Background agents are different",
      subtitle="Synchronous chat is forgiving. Background work isn't.",
      notes="Synchronous chat is forgiving. The user is staring at the screen, retries are free, a five "
            "minute Lambda is fine. Background agents are a different shape. They run for hours. They "
            "survive deploys. They pause for a human approval that won't arrive until tomorrow. And they "
            "have to wake up at the right step.",
      blocks=[
          {"t": "box", "title": "Chat agent vs background agent", "color": MAGENTA, "size": 13, "kvw": 3.3, "lines": [
              [seg("", DIM, 6), seg("                    ", DIM, 12, font=MONO),
               seg("chat", DIM, 13, font=MONO), seg("                  ", DIM, 12, font=MONO),
               seg("background", DIM, 13, font=MONO)],
              [seg("runtime         ", DIM, 12.5, font=MONO), seg("seconds", GREEN, 12.5, font=MONO), seg("                hours · days", YELLOW, 12.5, font=MONO)],
              [seg("user            ", DIM, 12.5, font=MONO), seg("staring at screen", GREEN, 12.5, font=MONO), seg("      offline", YELLOW, 12.5, font=MONO)],
              [seg("approval        ", DIM, 12.5, font=MONO), seg("immediate", GREEN, 12.5, font=MONO), seg("              tomorrow morning", YELLOW, 12.5, font=MONO)],
              [seg("crash           ", DIM, 12.5, font=MONO), seg("page refresh", GREEN, 12.5, font=MONO), seg("           lost work", YELLOW, 12.5, font=MONO)],
              [seg("deploy          ", DIM, 12.5, font=MONO), seg("reconnect", GREEN, 12.5, font=MONO), seg("              interrupted mid-task", YELLOW, 12.5, font=MONO)],
              [seg("observability   ", DIM, 12.5, font=MONO), seg("console", GREEN, 12.5, font=MONO), seg("                ???", YELLOW, 12.5, font=MONO)],
          ]},
      ])

# 04 — THE NAIVE FIX DOESN'T WORK ---------------------------------------------
slide(tag="ACT I · THE PROBLEM", accent=RED, title="The naive fix doesn't work",
      subtitle="A queue plus a database is 60% of an orchestrator, badly.",
      notes="You can build durable background agents with a queue and a database. But you'll reinvent "
            "sixty percent of what an honest durable execution layer already does, and you'll do it more "
            "poorly. Retries. Heartbeats. Resume from the right step. Approval suspension. Observability. "
            "None of these are application code you want to write.",
      blocks=[
          {"t": "box", "title": "What you'll re-implement (and get wrong)", "color": RED, "size": 13, "lines": [
              [seg("✗  ", RED, 13, bold=True, font=MONO), dim("durable step state machine", 13)],
              [seg("✗  ", RED, 13, bold=True, font=MONO), dim("heartbeat + stale-claim recovery", 13)],
              [seg("✗  ", RED, 13, bold=True, font=MONO), dim("retry policies with backoff", 13)],
              [seg("✗  ", RED, 13, bold=True, font=MONO), dim("suspension on approval / signal / event", 13)],
              [seg("✗  ", RED, 13, bold=True, font=MONO), dim("resume-at-the-right-step semantics", 13)],
              [seg("✗  ", RED, 13, bold=True, font=MONO), dim("cancellation propagation", 13)],
              [seg("✗  ", RED, 13, bold=True, font=MONO), dim("per-step / per-graph sandboxing", 13)],
              [seg("✗  ", RED, 13, bold=True, font=MONO), dim("structured observability (not just logs)", 13)],
              "",
              [seg("This isn't infrastructure you sprinkle on later. It's the substrate.", YELLOW, 13, bold=True, font=BODY)],
          ]},
      ])

# 05 — THE LAYER THAT DOESN'T CHANGE ------------------------------------------
slide(tag="ACT II · THE SHAPE OF THE ANSWER", accent=GREEN, title="The layer that doesn't change",
      subtitle="Underneath every named topology is a layer with a different velocity.",
      notes="Here's the thesis. Underneath every named topology — chains, ReAct, crews, swarms, "
            "background agents — there's a layer that doesn't change. Steps. State. Events. Retries. "
            "Observability. Smithers exists to be that layer.",
      blocks=[
          {"t": "layers"},
          {"t": "para", "align": PP_ALIGN.CENTER, "size": 13.5, "lines": [
              [seg("We are the bottom layer — the part you build on once and never throw away.", DIM, 13.5, font=BODY)],
          ]},
      ])

# 06 — THE FIVE PRIMITIVES ----------------------------------------------------
slide(tag="ACT II · THE SHAPE OF THE ANSWER", accent=GREEN, title="The five primitives",
      subtitle="Five capabilities the substrate has to provide — as uniform Effect-ts effects.",
      notes="Five things show up underneath every pattern. Durable steps. Persistent state. Parallel "
            "work. Event-driven control flow. Structured observability. Get these five right and the "
            "topology layer above becomes composable, not a runtime opinion.",
      blocks=[
          {"t": "box", "title": "Five primitives — uniform Effect-ts effects", "color": CYAN, "size": 13, "kvw": 0.5, "lines": [
              [seg("1  ", CYAN, 13, bold=True, font=MONO), seg("Durable steps      ", WHITE, 13, bold=True, font=HEAD), dim("<Task> · output decoded against a Zod schema, persisted", 12)],
              [seg("2  ", CYAN, 13, bold=True, font=MONO), seg("Persistent state   ", WHITE, 13, bold=True, font=HEAD), dim("every output schema becomes a typed SQLite table", 12)],
              [seg("3  ", CYAN, 13, bold=True, font=MONO), seg("Parallel work      ", WHITE, 13, bold=True, font=HEAD), dim("<Parallel> · structured concurrency on Effect fibers", 12)],
              [seg("4  ", CYAN, 13, bold=True, font=MONO), seg("Event-driven flow  ", WHITE, 13, bold=True, font=HEAD), dim("<Signal> <WaitForEvent> <Approval> · durable suspension", 12)],
              [seg("5  ", CYAN, 13, bold=True, font=MONO), seg("Observability      ", WHITE, 13, bold=True, font=HEAD), dim("Prometheus + SQLite event log · every transition is a row", 12)],
              "",
              [seg("A retry policy is a Schedule. A dependency is a Layer. A timeout is Effect.timeout.", DIM, 12.5, italic=True, font=BODY)],
              [seg("We didn't invent a parallel ecosystem.", DIM, 12.5, italic=True, font=BODY)],
          ]},
      ])

# 07 — THE FRAMEWORK TRAP -----------------------------------------------------
slide(tag="ACT II · THE SHAPE OF THE ANSWER", accent=GREEN, title="The framework trap",
      subtitle="Abstract the primitives. Not the topology.",
      notes="Agent frameworks aren't libraries. They're bets on which agent pattern wins. When the "
            "pattern shifts, you don't refactor. You rewrite. Smithers doesn't pick a topology for you. "
            "It hands you a primitive — a durable, retryable, observable task — and lets you "
            "compose whatever shape your problem needs.",
      blocks=[
          {"t": "box", "title": "Frameworks that age out vs frameworks that don't", "color": GREEN, "size": 13, "kvw": 5.6, "lines": [
              [seg("Topology-shaped", RED, 13, bold=True, font=MONO), seg("                     Substrate-shaped", GREEN, 13, bold=True, font=MONO)],
              "",
              [seg("AutoGPT-style agent loops", DIM, 12.5, font=MONO), seg("            Temporal, Durable Functions", DIM, 12.5, font=MONO)],
              [seg("Crew / swarm runtimes", DIM, 12.5, font=MONO), seg("                Effect-ts schedules", DIM, 12.5, font=MONO)],
              [seg("Graph-shaped DSL frameworks", DIM, 12.5, font=MONO), seg("          Smithers", DIM, 12.5, font=MONO)],
              "",
              [seg("A framework that abstracts the substrate ages fine.", DIM, 12.5, font=BODY)],
              [seg("A framework that abstracts the topology ages out when the topology does.", DIM, 12.5, font=BODY)],
              [seg("The mistake is conflating the two and throwing both away when one expires.", TEXT, 12.5, font=BODY)],
          ]},
      ])

# 08 — ONE-SIZE-FITS-ALL ORCHESTRATORS ----------------------------------------
slide(tag="ACT II · THE SHAPE OF THE ANSWER", accent=GREEN, title="One-size-fits-all orchestrators",
      subtitle="Gstack. Paperclips. Smithers is one layer down.",
      notes="Tools like Gstack and Paperclips are one-size-fits-all orchestrators. They're useful, "
            "they're real products, and they're also topology-shaped — they bet on what the right "
            "agent pipeline looks like, and they ship that pipeline. Smithers is one layer down from "
            "that. We believe the correct level of abstraction isn't a pre-built orchestrator. It's a "
            "framework you use to build your own, custom-fitted to your problem. Your taste, your "
            "topology, your shape. We took Gstack, an existing high-token agentic workflow, and cut it by "
            "roughly eighty percent of its lines of code just by composing Smithers components instead of "
            "hand-writing the orchestration.",
      blocks=[
          {"t": "box", "title": "Two layers of abstraction. Pick the right one.", "color": GREEN, "size": 12.5, "lines": [
              [seg("ONE-SIZE-FITS-ALL ORCHESTRATOR", RED, 12.5, bold=True, font=MONO), dim("     what Gstack / Paperclips / etc. ship", 11.5)],
              [seg("└── opinionated pipeline (planner → coder → reviewer → land)", DIM, 12, font=MONO)],
              [seg("    ├── prompts you can't easily change", MUTED, 12, font=MONO)],
              [seg("    ├── topology you can't easily reshape", MUTED, 12, font=MONO)],
              [seg("    └── ages out when the meta shifts", MUTED, 12, font=MONO)],
              "",
              [seg("ORCHESTRATION FRAMEWORK", GREEN, 12.5, bold=True, font=MONO), dim("            what Smithers ships", 11.5)],
              [seg("└── primitives (durable steps, retries, suspension, state)", DIM, 12, font=MONO)],
              [seg("    └── you compose your own orchestrator on top", MUTED, 12, font=MONO)],
          ]},
          {"t": "box", "color": YELLOW, "size": 12.5, "lines": [
              [seg("Gstack rewritten as Smithers components: ~80% fewer lines.", YELLOW, 13.5, bold=True, font=BODY)],
              [seg("Custom-fitted beats one-size-fits-all every time you can afford to write it — and with agents writing the orchestrator, you can always afford it.", DIM, 12, font=BODY)],
          ]},
      ])

# 09 — PATTERNS AS COMPONENTS -------------------------------------------------
slide(tag="ACT II · THE SHAPE OF THE ANSWER", accent=GREEN, title="Patterns as components",
      subtitle="We surveyed the field. Anything we saw twice became a component.",
      notes="We did deep research across every agentic orchestration framework and library we could find "
            "— LangGraph, Crew, Inngest, Temporal, AutoGen, Mastra, academic papers, vendor blog "
            "posts, open-source repos. Any pattern we saw more than once, and felt deserved promotion, we "
            "abstracted into a Smithers component. Anything that didn't quite earn a component lives in "
            "the examples folder as a recipe you can copy. Review loops. Optimizers. Scan-fix-verify. "
            "Panels. Debates. Escalation chains. Sagas. Every one ships as a composition on top of the "
            "substrate. None baked into the runtime. You can read the source. You can fork it. When the "
            "next pattern with no name yet shows up — and it will — you compose it from the same "
            "primitives.",
      blocks=[
          {"t": "box", "color": GREEN, "size": 12, "lines": [
              [seg("We surveyed every agentic orchestration framework we could find — vendors, OSS, papers — and codified what we saw repeatedly.", DIM, 12, font=BODY)],
              "",
              [seg("seen 2+ times · earned promotion", GREEN, 12, bold=True, font=MONO), seg("     →  built-in component", CYAN, 12, font=MONO)],
              [seg("seen, but project-specific", YELLOW, 12, bold=True, font=MONO), seg("           →  examples/ folder (101 files)", CYAN, 12, font=MONO)],
          ]},
          {"t": "box", "title": "Patterns shipped as components", "color": CYAN, "size": 11.5, "kvw": 2.7, "lines": [
              kv("<ReviewLoop>", "producer + reviewer, loop until approved", CYAN, 11.5),
              kv("<Optimizer>", "generator + evaluator, loop until score", CYAN, 11.5),
              kv("<ScanFixVerify>", "scanner → parallel fixers → verifier, retry survivors", CYAN, 11.5),
              kv("<Panel>", "N reviewers in parallel, moderator synthesizes", CYAN, 11.5),
              kv("<Debate>", "proposer vs opponent for N rounds, judge decides", CYAN, 11.5),
              kv("<GatherAndSynthesize>", "fan out, fan in", CYAN, 11.5),
              kv("<ClassifyAndRoute>", "classifier → category specialists in parallel", CYAN, 11.5),
              kv("<EscalationChain>", "tier 1 → tier 2 → human if confidence low", CYAN, 11.5),
              kv("<Supervisor>", "boss plans, workers execute, boss re-delegates", CYAN, 11.5),
              kv("<Saga>", "forward steps + compensations on failure", CYAN, 11.5),
          ]},
      ])

# 10 — THE RUNTIME LOOP -------------------------------------------------------
slide(tag="ACT III · HOW IT WORKS", accent=CYAN, title="The runtime loop",
      subtitle="Render → execute → persist → re-render.",
      notes="Mechanically, Smithers is a loop. Render your workflow tree. Extract the list of tasks. "
            "Execute the ones that are ready. Persist their outputs to SQLite. Re-render against the new "
            "state. That is the entire model.",
      blocks=[
          {"t": "runtime"},
      ])

# 11 — THE AUTHORING LAYER ----------------------------------------------------
slide(tag="ACT III · HOW IT WORKS", accent=CYAN, title="The authoring layer",
      subtitle="The fourth layer — legible to the agents that edit it.",
      notes="Inngest's three layer model is missing a fourth. In twenty twenty six a lot of workflow "
            "code is written and re-tuned by other agents. The authoring surface has to be legible to the "
            "agents that increasingly edit it, and to the humans auditing what those agents wrote.",
      blocks=[
          {"t": "box", "title": "Why TypeScript + JSX for authoring", "color": YELLOW, "size": 13, "lines": [
              [seg("TypeScript", WHITE, 13.5, bold=True, font=HEAD), seg("  because prompts are template strings.", TEXT, 13, font=BODY)],
              [seg("Interpolate. Refactor. Type-check. No DSL.", DIM, 12.5, font=BODY)],
              "",
              [seg("JSX", WHITE, 13.5, bold=True, font=HEAD), seg("  because agents are disproportionately good at it.", TEXT, 13, font=BODY)],
              [seg("React is the densest domain in any LLM's training corpus. Agents write it", DIM, 12.5, font=BODY)],
              [seg("fluently. Humans audit declarative trees better than imperative graphs.", DIM, 12.5, font=BODY)],
              "",
              [seg("MDX", WHITE, 13.5, bold=True, font=HEAD), seg("  because prompt fragments should compose like components.", TEXT, 13, font=BODY)],
              "",
              [seg("One bet among many — the substrate would work with any authoring surface. We picked the one agents already speak.", DIM, 12, italic=True, font=BODY)],
          ]},
      ])

# 12 — A WORKFLOW -------------------------------------------------------------
slide(tag="ACT III · HOW IT WORKS", accent=CYAN, title="A workflow",
      subtitle="Three tasks. One real conditional.",
      notes="Here is a real workflow. Three tasks in a sequence. Sequence already enforces order — "
            "fix waits for analyze without you doing anything. The interesting bit is the middle one: If "
            "analyze says a security review is required, an Approval mounts and the workflow durably "
            "suspends until a human answers. If it doesn't, the approval never exists. That conditional "
            "is real business logic flowing through J S X, and the entire shape of the run can change "
            "based on what the agent found.",
      blocks=[
          {"t": "split", "lw": 7.35, "file": "review.tsx", "code":
           r'''<Workflow name="review">
  <Sequence>
    <Task id="analyze" output={outputs.analysis}
          agent={analyst} retries={3}>
      Analyze {ctx.input.repo}@{ctx.input.sha}
    </Task>

    {analysis?.requiresSecurityReview && (
      <Approval
        id="security-review"
        output={outputs.securityReview}
        request={{
          title:    `Security review`,
          summary:  analysis.summary,
          severity: analysis.severity,
        }}
        onDeny="fail"
      />
    )}

    <Task id="fix" output={outputs.fix} agent={fixer}>
      Fix the issues:
      {analysis?.issues.map((i) =>
        `- [${i.severity}] ${i.description}`)}
    </Task>
  </Sequence>
</Workflow>''',
           "box": {"title": "What this gives you, for free", "color": GREEN, "size": 11.5, "lines": [
               [seg("✓  ", GREEN, 12, bold=True), seg("Sequence orders execution. ", TEXT, 11.5, font=BODY), seg("No wiring.", DIM, 11.5, font=BODY)],
               "",
               [seg("✓  ", GREEN, 12, bold=True), seg("analyze's output is decoded against a Zod schema, persisted.", TEXT, 11.5, font=BODY)],
               "",
               [seg("✓  ", GREEN, 12, bold=True), seg("Real conditional", WHITE, 11.5, bold=True, font=BODY), seg(" — the plan reshapes per run.", TEXT, 11.5, font=BODY)],
               "",
               [seg("✓  ", GREEN, 12, bold=True), seg("Approval durably suspends. Costs zero while waiting.", TEXT, 11.5, font=BODY)],
               "",
               [seg("✓  ", GREEN, 12, bold=True), seg("Retries default-on. You shouldn't write that loop.", TEXT, 11.5, font=BODY)],
               "",
               [seg("✓  ", GREEN, 12, bold=True), seg("Crash anywhere → resume from the last frame.", TEXT, 11.5, font=BODY)],
           ]}},
      ])

# 13 — CONTROL FLOW + HITL ----------------------------------------------------
slide(tag="ACT IV · COMPONENTS", accent=CYAN, title="Control flow + human-in-the-loop",
      subtitle="The core JSX surface.",
      notes="Everything is built on top of nine primitives. Workflow, Task, Sequence, Parallel, Branch, "
            "Loop for control flow. Approval, Signal, and Wait-for-event for durable human-in-the-loop "
            "suspension.",
      blocks=[
          {"t": "box", "title": "Control flow", "color": CYAN, "size": 12, "kvw": 1.7, "lines": [
              kv("<Workflow>", "root — names the run, owns the SQLite namespace", CYAN),
              kv("<Task>", "durable step · 3 modes: agent · compute · static", CYAN),
              kv("<Sequence>", "children execute in order", CYAN),
              kv("<Parallel>", "children execute concurrently, maxConcurrency knob", CYAN),
              kv("<Branch>", "if / then / else over persisted state", CYAN),
              kv("<Loop>", "until / maxIterations / onMaxReached", CYAN),
          ]},
          {"t": "box", "title": "Human-in-the-loop · durable suspension", "color": YELLOW, "size": 12, "kvw": 1.9, "lines": [
              kv("<Approval>", "pause for approve / deny, runtime exits, resumes when answered", YELLOW),
              kv("<HumanTask>", "structured ask — schema + form, durably waits for response", YELLOW),
              kv("<Signal>", "wake on  smithers signal <run> <name>", YELLOW),
              kv("<WaitForEvent>", "wake on a webhook, HTTP POST, or external trigger", YELLOW),
              "",
              [seg("A suspended run is a row, not a process. Costs zero while waiting.", DIM, 12, italic=True, font=BODY)],
          ]},
      ])

# 14 — ReviewLoop + Optimizer (usage) -----------------------------------------
slide(tag="ACT IV · COMPONENTS", accent=CYAN, title="<ReviewLoop> + <Optimizer>",
      subtitle="Producer/reviewer and generator/evaluator, shipped as JSX.",
      notes="Review-loop pairs a producer with a reviewer and loops until the reviewer approves. "
            "Optimizer pairs a generator with an evaluator and loops until a target score is reached. "
            "Both ship in the box. Both are forty lines of source you can read and copy.",
      blocks=[
          {"t": "code", "file": "usage.tsx", "text":
           r'''<ReviewLoop
  producer={coder}
  reviewer={[primaryReviewer, secondaryReviewer]}     // array = consensus
  produceOutput={outputs.code}
  reviewOutput={outputs.review}                       // must include approved: boolean
  maxIterations={5}
>
  Produce a function that {ctx.input.task}.
</ReviewLoop>

<Optimizer
  generator={promptEngineer}
  evaluator={evaluator}                               // agent or compute fn
  generateOutput={outputs.prompt}
  evaluateOutput={outputs.evaluation}                 // must include score: number
  targetScore={90}
  maxIterations={5}
>
  Generate a prompt for summarising legal documents.
</Optimizer>'''},
      ])

# 15 — ...AND HERE IS THEIR SOURCE (two columns) ------------------------------
slide(tag="ACT IV · COMPONENTS", accent=CYAN, title="…and here is their source",
      subtitle="Pattern components are just compositions. You can read them.",
      notes="And here is what those components actually are. Twenty lines of J S X each. Loop wrapping a "
            "Sequence wrapping two Tasks. Nothing baked into the runtime. You can read them. You can fork "
            "them. When the next pattern with no name yet shows up — and it will — you compose "
            "it from the same primitives.",
      blocks=[
          {"t": "code2", "size": 9.0, "lfile": "ReviewLoop.tsx", "rfile": "Optimizer.tsx",
           "left":
           r'''export function ReviewLoop({
  id = "review-loop", producer, reviewer,
  produceOutput, reviewOutput,
  maxIterations = 5,
  onMaxReached = "return-last",
  children,
}: ReviewLoopProps) {
  const reviewers = Array.isArray(reviewer)
    ? reviewer : [reviewer];
  return (
    <Loop id={id} until={false}
          maxIterations={maxIterations}
          onMaxReached={onMaxReached}>
      <Sequence>
        <Task id={`${id}-produce`}
              output={produceOutput}
              agent={producer}>
          {children}
        </Task>
        <Task id={`${id}-review`}
              output={reviewOutput}
              agent={reviewers}
              needs={{ produced: `${id}-produce` }}>
          Review and decide to approve.
        </Task>
      </Sequence>
    </Loop>
  );
}''',
           "right":
           r'''export function Optimizer({
  id = "optimizer", generator, evaluator,
  generateOutput, evaluateOutput,
  maxIterations = 10,
  onMaxReached = "return-last",
  children,
}: OptimizerProps) {
  const isAgent = typeof evaluator !== "function";
  return (
    <Loop id={id} until={false}
          maxIterations={maxIterations}
          onMaxReached={onMaxReached}>
      <Sequence>
        <Task id={`${id}-generate`}
              output={generateOutput}
              agent={generator}>
          {children}
        </Task>
        <Task id={`${id}-evaluate`}
              output={evaluateOutput}
              agent={isAgent ? evaluator : undefined}
              needs={{ cand: `${id}-generate` }}>
          {isAgent ? "Evaluate + score." : evaluator}
        </Task>
      </Sequence>
    </Loop>
  );
}'''},
      ])

# 16 — ScanFixVerify + Debate + Panel -----------------------------------------
slide(tag="ACT IV · COMPONENTS", accent=CYAN, title="<ScanFixVerify> + <Debate> + <Panel>",
      subtitle="Composable adversarial and parallel-fan-out patterns.",
      notes="Scan-fix-verify: a scanner finds issues, fixers run in parallel, a verifier confirms each "
            "fix, survivors retry. Debate: a proposer and an opponent argue for N rounds, a judge "
            "synthesises. Panel: N specialist reviewers in parallel, a moderator synthesises by vote, "
            "consensus, or merge.",
      blocks=[
          {"t": "code", "file": "patterns.tsx", "text":
           r'''<ScanFixVerify
  scanner={lintAgent}
  fixer={[fixerA, fixerB, fixerC]}                    // array cycles across issues
  verifier={verifyAgent}
  scanOutput={outputs.scan}     fixOutput={outputs.fix}
  verifyOutput={outputs.verify} reportOutput={outputs.report}
  maxConcurrency={4}            maxRetries={3}
/>

<Debate proposer={pro} opponent={con} judge={judge}
  rounds={3} verdictOutput={outputs.verdict}>
  Should we migrate from Postgres to ClickHouse?
</Debate>

<Panel reviewers={[security, perf, ux, infra]}
  moderator={pm} synthesis="consensus"
  reviewOutput={outputs.reviews} verdictOutput={outputs.verdict}>
  Review the RFC.
</Panel>'''},
      ])

# 17 — MORE PATTERNS ----------------------------------------------------------
slide(tag="ACT IV · COMPONENTS", accent=CYAN, title="More patterns",
      subtitle="All composable. None baked into the runtime.",
      notes="Supervisor is a boss agent that plans, dispatches to workers in parallel, reviews their "
            "work, and re-delegates failures. Saga runs forward steps with compensations that fire in "
            "reverse on failure. Kanban runs work items through a configurable column pipeline. "
            "Escalation chain tries tier one, escalates to tier two if confidence is low, then to a "
            "human. Classify-and-route sorts items into categories and dispatches to category "
            "specialists. All of these are JSX components on top of the substrate.",
      blocks=[
          {"t": "box", "title": "Built-in pattern components", "color": CYAN, "size": 11, "kvw": 2.6, "lines": [
              kv("<Supervisor>", "boss plans · workers parallel · re-delegate failures", CYAN, 11),
              kv("<Saga>", "forward steps + compensations on failure", CYAN, 11),
              kv("<Kanban>", "items flow through configurable columns", CYAN, 11),
              kv("<MergeQueue>", "serialise risky ops · single in-flight rule", CYAN, 11),
              kv("<EscalationChain>", "tier 1 → tier 2 → human on low confidence", CYAN, 11),
              kv("<ClassifyAndRoute>", "classifier → category specialists in parallel", CYAN, 11),
              kv("<GatherAndSynthesize>", "fan out to N sources · synthesise the result", CYAN, 11),
              kv("<CheckSuite>", "declare must-pass checks · gate downstream", CYAN, 11),
              kv("<DecisionTable>", "rule-based dispatch, deterministic", CYAN, 11),
              kv("<Poller>", "poll external condition with backoff", CYAN, 11),
              kv("<Runbook>", "declarative ops procedures", CYAN, 11),
              kv("<DriftDetector>", "detect changes against a baseline", CYAN, 11),
              kv("<ContentPipeline>", "staged content transforms", CYAN, 11),
              kv("<TryCatchFinally>", "structured error handling · <Timer> durable sleep", CYAN, 11),
          ]},
      ])

# 18 — ISOLATION --------------------------------------------------------------
slide(tag="ACT IV · COMPONENTS", accent=CYAN, title="Isolation",
      subtitle="Sandboxes, worktrees, subflows, and sub-workflows.",
      notes="Tasks need isolation. Sometimes per-step, sometimes per-graph. Sandbox runs a child "
            "workflow or a single step in an isolated runtime. Worktree gives each parallel agent its own "
            "git worktree so they don't fight over port five thousand one seventy three. Subflow embeds "
            "another workflow file as a single node. Super-smithers spawns a whole nested workflow with "
            "its own database scope.",
      blocks=[
          {"t": "code", "file": "isolation.tsx", "text":
           r'''const remoteVmProvider = {
  id: "remote-vm",
  async run(request) { return runRemoteVm(request); },
};

<Worktree path=".worktrees/feature-a" baseBranch="main">
  <Parallel>
    <Task id="fix-a" output={outputs.patch} agent={fixer}>Fix issue A</Task>
    <Task id="fix-b" output={outputs.patch} agent={fixer}>Fix issue B</Task>
  </Parallel>
</Worktree>

<Sandbox id="exec" provider={remoteVmProvider}
  workflow={testWorkflow} input={{ patch: outputs.patch }}
  output={outputs.sandbox} />

<Subflow workflow={reviewWorkflow} input={{ repo, sha }} output={outputs.review} />
<SuperSmithers strategy={strategyDoc} agent={engineer} reportOutput={outputs.report} />'''},
          {"t": "para", "align": PP_ALIGN.CENTER, "size": 12, "lines": [
              [seg("Pluggable sandbox providers: VM adapters, local transports, or custom runners.", DIM, 12, italic=True, font=BODY)],
          ]},
      ])

# 19 — Aspects ----------------------------------------------------------------
slide(tag="ACT V · BEYOND THE JSX SURFACE", accent=MAGENTA, title="<Aspects>",
      subtitle="Cross-cutting budgets — tokens, latency, cost.",
      notes="Wrap any subtree in Aspects to propagate budgets to descendant tasks. Token budget. Latency "
            "S L O. Cost budget. Each can fail, warn, or skip-remaining when exceeded. Nested Aspects "
            "inherit. Inner fields override per-config.",
      blocks=[
          {"t": "code", "file": "budgeted.tsx", "text":
           r'''<Workflow name="budgeted">
  <Aspects
    tokenBudget={{ max: 100_000, perTask: 25_000, onExceeded: "warn" }}
    latencySlo ={{ maxMs: 30_000,                onExceeded: "fail" }}
    costBudget ={{ maxUsd: 5.0,                  onExceeded: "skip-remaining" }}
    tracking   ={{ tokens: true, latency: true, cost: true }}
  >
    <Task id="analyse" output={outputs.analysis} agent={analyst}>...</Task>
    <Task id="review"  output={outputs.review}   agent={reviewer}>...</Task>
  </Aspects>
</Workflow>'''},
          {"t": "box", "color": GREEN, "size": 12, "lines": [
              [seg("Token / latency / cost counts accumulate per-run, reset on resume, and emit to Prometheus.", DIM, 12, font=BODY)],
              [seg("A retry policy is a Schedule. A timeout is Effect.timeout. You're composing Effect-ts primitives, not bespoke knobs.", DIM, 12, font=BODY)],
          ]},
      ])

# 20 — SCORERS / EVALS --------------------------------------------------------
slide(tag="ACT V · BEYOND THE JSX SURFACE", accent=MAGENTA, title="Scorers · evals",
      subtitle="Quality gates on every task output.",
      notes="Every task can have scorers attached. Smithers ships built-ins for faithfulness, relevancy, "
            "schema adherence, toxicity, and latency. Plus an L L M judge scorer for everything else. "
            "Scores aggregate. They emit to Prometheus. They show up in the dashboard, and you can query "
            "them from the C L I with smithers scores.",
      blocks=[
          {"t": "code", "file": "scored.tsx", "text":
           r'''import { llmJudge, faithfulness, relevancy, schemaAdherence } from
  "smithers-orchestrator/scorers";

<Task
  id="answer"
  output={outputs.answer}
  agent={researcher}
  scorers={[
    faithfulness({ source: "context" }),               // grounded in the input
    relevancy({ question: ctx.input.question }),       // answers the question
    schemaAdherence(),                                 // matches the Zod schema
    llmJudge({                                         // custom LLM judge
      model: anthropic("claude-sonnet-4-20250514"),
      rubric: "Score 0-100 on clarity and concision.",
    }),
  ]}
>
  {`Answer: ${ctx.input.question}`}
</Task>'''},
          {"t": "para", "align": PP_ALIGN.CENTER, "size": 12, "lines": [
              [seg("smithers scores <run-id>", CYAN, 12, bold=True, font=MONO),
               seg("    — tabular view, per-task, per-scorer", DIM, 12, font=BODY)],
          ]},
      ])

# 21 — MEMORY -----------------------------------------------------------------
slide(tag="ACT V · BEYOND THE JSX SURFACE", accent=MAGENTA, title="Memory",
      subtitle="Cross-run state. Outputs are per-run, memory survives.",
      notes="Outputs are per-run. Memory is per-namespace and survives every workflow execution. Three "
            "layers — facts with optional T T L, ordered message history, and maintenance for "
            "compaction. On any task you can set memory dot recall to auto-inject the top K most relevant "
            "past facts into the prompt.",
      blocks=[
          {"t": "code", "file": "memory.tsx", "text":
           r'''import { createMemoryStore } from "smithers-orchestrator/memory";
const ns = { kind: "workflow", id: "code-review" };

<Task
  id="review"
  output={outputs.review}
  agent={reviewer}
  memory={{
    recall: { namespace: ns, topK: 3 },              // auto-inject past facts
    save:   { namespace: ns,
              key: ({ run }) => `review:${run.id}` }, // persist this output
  }}
>
  Review {ctx.input.diff}
</Task>

// Or imperatively:
store.setFact(ns, "code-style", { tabs: 2, semi: true }, 30 * 24 * 3600_000);'''},
      ])

# 22 — TOOLS / AGENTS / MCP ---------------------------------------------------
slide(tag="ACT V · BEYOND THE JSX SURFACE", accent=MAGENTA, title="Tools · agents · MCP",
      subtitle="Sandboxed tools, any agent CLI, OpenAPI generator, MCP server.",
      notes="Smithers ships read, write, edit, bash, and grep tools with path containment — pass "
            "--root to set the sandbox boundary. Agents are pluggable: claude, codex, antigravity, kimi, "
            "amp, forge, or anything that implements the Agent interface. Agent fallback lets you write "
            "agent equals an array — primary first, fallback on failure. Smithers openapi generates "
            "AI SDK tools from an OpenAPI spec. And smithers itself can run as an M C P server with "
            "smithers mcp add.",
      blocks=[
          {"t": "box", "title": "Tools, agents, integrations", "color": CYAN, "size": 12, "kvw": 1.9, "lines": [
              kv("Built-in tools", "read · write · edit · bash · grep · ls (path-contained)", CYAN),
              kv("Agent fallback", "agent={[claude, codex]}  // claude first, codex on fail", CYAN),
              kv("Agent runtimes", "claude · codex · antigravity · kimi · amp · forge · Effect-native", CYAN),
              kv("MDX prompts", "prompt fragments with typed props · imports compose", CYAN),
              kv("OpenAPI tools", "smithers openapi <spec> → typed AI SDK tool surface", CYAN),
              kv("MCP server", "smithers mcp add  // call workflows from any MCP agent", CYAN),
              kv("Skills sync", "smithers skills add  // bundle skills into agent dirs", CYAN),
          ]},
      ])

# 23 — HOT MODE ---------------------------------------------------------------
slide(tag="ACT V · BEYOND THE JSX SURFACE", accent=MAGENTA, title="Hot mode",
      subtitle="Edit prompts and workflow code while a run is in flight.",
      notes="Pass --hot true to smithers up. Edit the workflow file, edit any M D X prompt, save. The "
            "runtime detects the change, re-renders the tree against the current persisted state, and "
            "continues from where it was. Tasks that already completed stay in the database. Tasks that "
            "haven't run yet use the new code.",
      blocks=[
          {"t": "box", "title": "Hot reload in practice", "color": MAGENTA, "size": 12.5, "lines": [
              [seg("$ ", MUTED, 12.5, font=MONO), seg("smithers up workflow.tsx --hot true", WHITE, 12.5, bold=True, font=MONO)],
              [seg("[00:00:02] ✓ analyze (attempt 1)", DIM, 12.5, font=MONO)],
              [seg("[00:00:02] → fix (attempt 1, iteration 0)", DIM, 12.5, font=MONO)],
              "",
              [seg("      …meanwhile, in another pane:", MUTED, 12.5, font=MONO)],
              [seg("vim .smithers/prompts/fix.mdx", GREEN, 12.5, font=MONO)],
              "",
              [seg("[00:00:14] ↻ hot reload — re-rendering tree", DIM, 12.5, font=MONO)],
              [seg("[00:00:14] → fix (attempt 2, iteration 0)  ← new prompt, same DB", DIM, 12.5, font=MONO)],
              "",
              [seg("Frame numbers march on. The DB is your time machine.", YELLOW, 12.5, bold=True, font=BODY)],
          ]},
      ])

# 24 — CLI RUN LIFECYCLE ------------------------------------------------------
slide(tag="ACT VI · THE CLI", accent=YELLOW, title="CLI · run lifecycle",
      subtitle="Compose-style commands for managing runs.",
      notes="The C L I follows compose semantics. Up. Ps. Inspect. Logs. Cancel. Down. Plus supervise, "
            "which polls for stale heartbeats and auto-resumes orphaned runs.",
      blocks=[
          {"t": "box", "title": "Lifecycle commands", "color": CYAN, "size": 12, "lines": [
              [seg("smithers up", CYAN, 12, bold=True, font=MONO), seg(" <file>", TEXT, 12, font=MONO), seg("         start a run · -d detached · --serve for HTTP API", DIM, 11.5, font=BODY)],
              [seg("smithers ps", CYAN, 12, bold=True, font=MONO), seg("                 active · paused · recently completed", DIM, 11.5, font=BODY)],
              [seg("smithers inspect", CYAN, 12, bold=True, font=MONO), seg(" <run>", TEXT, 12, font=MONO), seg("      structured state · -w for watch mode", DIM, 11.5, font=BODY)],
              [seg("smithers logs", CYAN, 12, bold=True, font=MONO), seg(" <run>", TEXT, 12, font=MONO), seg("         NDJSON event log · streamable", DIM, 11.5, font=BODY)],
              [seg("smithers node", CYAN, 12, bold=True, font=MONO), seg(" <run> <node>", TEXT, 12, font=MONO), seg("  per-task detail: attempts, tool calls, output", DIM, 11.5, font=BODY)],
              [seg("smithers cancel", CYAN, 12, bold=True, font=MONO), seg(" <run>", TEXT, 12, font=MONO), seg("       safely halt agents · terminate", DIM, 11.5, font=BODY)],
              [seg("smithers down", CYAN, 12, bold=True, font=MONO), seg("               cancel ALL active runs · compose-down energy", DIM, 11.5, font=BODY)],
              [seg("smithers supervise", CYAN, 12, bold=True, font=MONO), seg("          auto-resume stale runs · --stale-threshold 30s", DIM, 11.5, font=BODY)],
              [seg("smithers why", CYAN, 12, bold=True, font=MONO), seg(" <run>", TEXT, 12, font=MONO), seg("          explain why a run is blocked / paused", DIM, 11.5, font=BODY)],
              "",
              [seg("--format toon | json | yaml | md | jsonl  — pick your output", DIM, 11.5, italic=True, font=MONO)],
          ]},
      ])

# 25 — CLI TIME TRAVEL --------------------------------------------------------
slide(tag="ACT VI · THE CLI", accent=YELLOW, title="CLI · time travel",
      subtitle="Every frame is a row. Replay, fork, diff.",
      notes="Every render frame is a database row. Timeline lists them. Fork branches from any frame. "
            "Replay forks and resumes. Diff compares two snapshots. Rewind walks a run back. Timetravel "
            "reverts the filesystem alongside the database. Retry-task re-runs one node without resetting "
            "the rest.",
      blocks=[
          {"t": "box", "title": "Time-travel commands", "color": MAGENTA, "size": 12, "lines": [
              [seg("smithers timeline", MAGENTA, 12, bold=True, font=MONO), seg(" <run>", TEXT, 12, font=MONO), seg("            list every frame · --tree includes forks", DIM, 11.5, font=BODY)],
              [seg("smithers fork", MAGENTA, 12, bold=True, font=MONO), seg("     <run> --frame N", TEXT, 12, font=MONO), seg("  branch from any frame", DIM, 11.5, font=BODY)],
              [seg("smithers replay", MAGENTA, 12, bold=True, font=MONO), seg("   <run> --frame N", TEXT, 12, font=MONO), seg("  fork + immediately resume", DIM, 11.5, font=BODY)],
              [seg("smithers rewind", MAGENTA, 12, bold=True, font=MONO), seg("   <run> --frame N", TEXT, 12, font=MONO), seg("  rewind in-place", DIM, 11.5, font=BODY)],
              [seg("smithers diff", MAGENTA, 12, bold=True, font=MONO), seg("     <a> <b>", TEXT, 12, font=MONO), seg("          DiffBundle as a unified diff", DIM, 11.5, font=BODY)],
              [seg("smithers timetravel", MAGENTA, 12, bold=True, font=MONO), seg(" <wf> -r <run>", TEXT, 12, font=MONO), seg("  + revert filesystem state", DIM, 11.5, font=BODY)],
              [seg("smithers retry-task", MAGENTA, 12, bold=True, font=MONO), seg(" <run> -n <id>", TEXT, 12, font=MONO), seg("  re-run one node, resume the workflow", DIM, 11.5, font=BODY)],
              [seg("smithers revert", MAGENTA, 12, bold=True, font=MONO), seg("   <run> --attempt N", TEXT, 12, font=MONO), seg(" revert FS to a prior task attempt", DIM, 11.5, font=BODY)],
              "",
              [seg("State lives in SQLite. Time travel is just SELECT, fork, REPLACE.", DIM, 11.5, italic=True, font=BODY)],
          ]},
      ])

# 26 — CLI CRON / ALERTS / HITL -----------------------------------------------
slide(tag="ACT VI · THE CLI", accent=YELLOW, title="CLI · cron · alerts · HITL",
      subtitle="Durable scheduling, alerts, signals, human queues.",
      notes="Cron schedules a recurring workflow trigger durably — survives restarts. Alerts is the "
            "durable equivalent of a pager — a run can raise one and humans can resolve it. Signal "
            "wakes a workflow blocked on wait-for-event. Approve, deny, and human resolve the "
            "human-in-the-loop suspension points. Memory exposes cross-run facts. Events queries the "
            "structured event log. Token mints short-lived gateway bearer tokens.",
      blocks=[
          {"t": "box", "title": "Durable extras", "color": YELLOW, "size": 12, "kvw": 2.5, "lines": [
              kv("smithers cron", "schedule recurring triggers · durable", YELLOW),
              kv("smithers alerts", "list / resolve durable alert instances", YELLOW),
              kv("smithers signal", "wake a run blocked on <WaitForEvent>", YELLOW),
              kv("smithers approve / deny", "resolve <Approval> gates", YELLOW),
              kv("smithers human", "list / resolve <HumanTask> requests", YELLOW),
              kv("smithers memory", "query cross-run facts · semantic recall", YELLOW),
              kv("smithers events", "query event log · NDJSON · filters · grouping", YELLOW),
              kv("smithers token", "issue / revoke short-lived Gateway bearer tokens", YELLOW),
              kv("smithers scores", "view scorer results for a run", YELLOW),
              kv("smithers ask", "ask the docs MCP server in natural language", YELLOW),
          ]},
      ])

# 27 — LIVE DEMO: DURABILITY --------------------------------------------------
slide(tag="ACT VII · LIVE DEMO", accent=GREEN, title="Live demo — durability",
      subtitle="Real workflow. Real crash. Real resume.   (run live in the talk)",
      notes="I'm going to run a real workflow now. Three tasks. I'll kill the process while the second "
            "task is running, then resume it. Watch what gets re-run and what doesn't.",
      blocks=[
          {"t": "split", "lw": 6.7, "file": "sample.tsx", "code":
           r'''<Workflow name="ship-it">
  <Sequence>
    <Task id="research" output={outputs.research}>
      {/* fast */}
      {async () =>
        ({ message: "scanned repo" })}
    </Task>
    <Task id="plan" output={outputs.plan}>
      {/* slow on purpose */}
      {async () => { await wait(7);
        return { message: "drafted plan" }; }}
    </Task>
    <Task id="implement" output={outputs.implement}>
      {async () =>
        ({ message: "patches applied" })}
    </Task>
  </Sequence>
</Workflow>''',
           "box": {"title": "Kill it mid-run, then resume", "color": GREEN, "size": 11.5, "lines": [
               [seg("$ ", MUTED, 11.5, font=MONO), seg("smithers up sample.tsx", WHITE, 11.5, bold=True, font=MONO)],
               [seg("✓ research        → plan running…", DIM, 11.5, font=MONO)],
               [seg("^C", RED, 11.5, bold=True, font=MONO), seg("  (closing the laptop)", MUTED, 11.5, font=MONO)],
               "",
               [seg("$ ", MUTED, 11.5, font=MONO), seg("smithers up sample.tsx --resume true", WHITE, 11.5, bold=True, font=MONO)],
               "",
               [seg("✓ ", GREEN, 11.5, bold=True), seg("research skipped — already in the DB", DIM, 11.5, font=BODY)],
               [seg("✓ ", GREEN, 11.5, bold=True), seg("plan re-ran as attempt 2 (was interrupted)", DIM, 11.5, font=BODY)],
               [seg("✓ ", GREEN, 11.5, bold=True), seg("implement ran for the first time", DIM, 11.5, font=BODY)],
               "",
               [seg("no work lost.", GREEN, 12.5, bold=True, font=BODY)],
           ]}},
      ])

# 28 — LIVE DEMO: TIME TRAVEL -------------------------------------------------
slide(tag="ACT VII · LIVE DEMO", accent=GREEN, title="Live demo — time travel",
      subtitle="Every frame is a snapshot.   (run live in the talk)",
      notes="Every render of the workflow tree is committed to the database as a frame. Let me show you "
            "the frames from the run we just finished. Not a log of what happened — the actual state.",
      blocks=[
          {"t": "box", "color": MAGENTA, "size": 12, "lines": [
              [seg("$ ", MUTED, 12, font=MONO), seg("smithers timeline demo-1729…", WHITE, 12, bold=True, font=MONO)],
              [seg("  frame 0   render      research · plan · implement  (pending)", DIM, 11.5, font=MONO)],
              [seg("  frame 1   execute     research ✓", DIM, 11.5, font=MONO)],
              [seg("  frame 2   execute     plan (attempt 1, interrupted)", DIM, 11.5, font=MONO)],
              [seg("  frame 3   resume      plan (attempt 2) ✓", DIM, 11.5, font=MONO)],
              [seg("  frame 4   execute     implement ✓  ·  run complete", DIM, 11.5, font=MONO)],
          ]},
          {"t": "box", "title": "What you can do with frames", "color": MAGENTA, "size": 12, "kvw": 2.2, "lines": [
              kv("smithers fork", "branch from any frame", WHITE),
              kv("smithers replay", "re-execute from a checkpoint", WHITE),
              kv("smithers diff", "compare two snapshots", WHITE),
              kv("smithers timetravel", "rewind FS state, edit an output, replay forward", WHITE),
              "",
              [seg("Git history for AI workflows. The actual state, not just logs of it.", DIM, 12, italic=True, font=BODY)],
          ]},
      ])

# 29 — OBSERVABILITY ----------------------------------------------------------
slide(tag="ACT VIII · PRODUCTION", accent=CYAN, title="Observability",
      subtitle="Grafana, Prometheus, Tempo, OTLP — one command.",
      notes="Smithers observability brings up a full local stack with one command. Grafana for "
            "dashboards. Prometheus for metrics. Tempo for traces. An OTLP collector so any task can emit "
            "spans. Every state transition, every attempt, every retry is already a row in the event log "
            "— you can SQL it.",
      blocks=[
          {"t": "box", "title": "Bring up the full local stack", "color": GREEN, "size": 12.5, "lines": [
              [seg("$ ", MUTED, 12.5, font=MONO), seg("smithers observability up", WHITE, 12.5, bold=True, font=MONO)],
              [seg("→ Grafana       http://localhost:3000", DIM, 12, font=MONO)],
              [seg("→ Prometheus    http://localhost:9090", DIM, 12, font=MONO)],
              [seg("→ Tempo         http://localhost:3200", DIM, 12, font=MONO)],
              [seg("→ OTLP collector :4317 (gRPC) :4318 (HTTP)", DIM, 12, font=MONO)],
              "",
              [seg("Pre-wired dashboards. No setup. Just `up` and open the link.", DIM, 12, font=BODY)],
              "",
              [seg("smithers up workflow.tsx --serve --metrics", GREEN, 12, bold=True, font=MONO)],
              [seg("→ HTTP API at :7331  /v1/runs  /v1/runs/:id/events (SSE)", DIM, 12, font=MONO)],
              [seg("→ Prometheus scrape endpoint at :7331/metrics", DIM, 12, font=MONO)],
          ]},
      ])

# 30 — HTTP SERVER / WORKFLOW APPS --------------------------------------------
slide(tag="ACT VIII · PRODUCTION", accent=CYAN, title="HTTP server · workflow apps",
      subtitle="A workflow can serve its own React UI.",
      notes="Smithers up dash dash serve runs an H T T P server alongside the workflow. Routes for "
            "listing runs, inspecting one, and a server-sent-events stream for live updates. Bearer-token "
            "auth. Workflows can also serve their own front-end — a workflow ships an H T M L bundle "
            "alongside the T S X file, and Smithers hands it to any client that asks. Workflows as full "
            "apps, not just task graphs.",
      blocks=[
          {"t": "code", "file": "serve.sh", "text":
           r'''// .smithers/workflows/kanban.tsx
// .smithers/workflows/kanban.frontend/
//   index.html · assets/ · manifest.json   ← React app, served by smithers

bunx smithers-orchestrator up workflow.tsx \
  --serve --port 7331 \
  --auth-token "$SMITHERS_API_KEY" \
  --metrics

GET  /v1/runs                       # list
GET  /v1/runs/:id                   # detail
GET  /v1/runs/:id/events            # SSE stream
GET  /v1/workflows/:name/app/*      # served front-end
GET  /metrics                       # Prometheus'''},
      ])

# 31 — EFFECT-TS API ----------------------------------------------------------
slide(tag="ACT VIII · PRODUCTION", accent=CYAN, title="Effect-ts API",
      subtitle="The substrate, exposed.",
      notes="Underneath the J S X surface is Effect-ts. For users who already think in Effect dot gen, "
            "Smithers exposes a slightly lower-level Effect A P I with full access to schedules, layers, "
            "fibers, and resource lifetimes. Same substrate. Different authoring surface. You can mix both "
            "in one workflow.",
      blocks=[
          {"t": "code", "file": "effect-api.ts", "text":
           r'''import { Smithers } from "smithers-orchestrator";
import { Effect, Schema } from "effect";

const G = Smithers.workflow({
  name: "review",
  input: Schema.Struct({ repo: Schema.String, sha: Schema.String }),
});

const analyze = G.step("analyze", {
  output: Schema.Struct({ summary: Schema.String, risk: Schema.Literal("low","med","high") }),
  timeout: "2m",
  retry: { maxAttempts: 3, backoff: "exponential", initialDelay: "1s" },
  run: ({ input, heartbeat, signal }) =>
    Effect.gen(function* () {
      heartbeat({ phase: "analyzing" });
      return yield* analyzeRepo(input, { signal });
    }),
});'''},
          {"t": "para", "align": PP_ALIGN.CENTER, "size": 12, "lines": [
              [seg("A retry policy is a Schedule. A dependency is a Layer. A timeout is Effect.timeout.", DIM, 12, italic=True, font=BODY)],
          ]},
      ])

# 32 — 101 EXAMPLES -----------------------------------------------------------
slide(tag="ACT VIII · PRODUCTION", accent=CYAN, title="101 examples",
      subtitle="A starter zoo. Pick one and edit.",
      notes="The examples folder ships over a hundred real workflows. Review loops, debates, optimizers, "
            "parallel ticket processors, refactor pipelines, kanban boards, supervisors, classifier "
            "switchboards, alert suppressors, doc sync, repo janitors, ransomware isolation coordinators, "
            "financial inbox guards, and a Ralph loop that keeps going until the work is done. All on the "
            "same substrate.",
      blocks=[
          {"t": "box", "title": "Pick something to start from", "color": CYAN, "size": 12, "kvw": 2.4, "lines": [
              kv("code-review-loop", "producer + reviewer + agent fallback", CYAN),
              kv("debate", "two agents argue, a judge decides", CYAN),
              kv("supervisor", "boss plans, workers in worktrees, re-delegates", CYAN),
              kv("parallel-tickets", "ingest tickets → triage → fix in waves", CYAN),
              kv("kanban", "configurable column pipeline", CYAN),
              kv("prompt-optimizer", "generator + evaluator + target score", CYAN),
              kv("migration", "schema migration with checkpoints + revert", CYAN),
              kv("repo-janitor", "scan, fix, verify across a whole repo", CYAN),
              kv("friday-bot", "weekly digest, cron-triggered", CYAN),
              kv("ralph-loop", "keep going until done", CYAN),
              "",
              [seg("101 files in examples/. ", DIM, 12, font=BODY),
               seg("bunx smithers-orchestrator init", GREEN, 12, bold=True, font=MONO),
               seg("  scaffolds them into your repo.", DIM, 12, font=BODY)],
          ]},
      ])

# 33 — ONE MORE THING ---------------------------------------------------------
slide(tag="ACT IX · WHAT SMITHERS IS", accent=TEAL, title="One more thing",
      subtitle="This slideshow is itself a Smithers workflow.",
      notes="Everything you just watched — every slide, the live crash demo, the time travel — "
            "is itself a Smithers workflow. Each slide is rendered by a Task. The keyboard nav is wrapped "
            "in the workflow tag. If I had killed this process two slides ago and resumed, you would have "
            "picked up where I left off.",
      blocks=[
          {"t": "code", "file": "demo.tsx", "text":
           r'''export default smithers((ctx) => (
  <Workflow name="demo">
    <Task id="slideshow" output={outputs.slideshow}>
      {async () => {
        for (let i = ctx.input.startAt; i < SLIDES.length; i++) {
          await renderSlide(SLIDES[i], i, ctx);
          await waitForArrowKey();
        }
        return { finished: true };
      }}
    </Task>
  </Workflow>
));'''},
          {"t": "para", "align": PP_ALIGN.CENTER, "size": 12, "lines": [
              [seg(".smithers/workflows/demo.tsx — the whole deck is a single Task.", DIM, 12, italic=True, font=BODY)],
          ]},
      ])

# 34 — THREE LAYERS -----------------------------------------------------------
slide(tag="ACT IX · WHAT SMITHERS IS", accent=TEAL, title="Three layers",
      subtitle="smithers-orchestrator. The forge. The GUI.",
      notes="Smithers is three things. The orchestrator — what you just watched — shipped today "
            "on N P M. The forge — a J J native code host with cloud workspaces — in build, AGPL. "
            "And a native macOS app you can download right now.",
      blocks=[
          {"t": "box", "color": CYAN, "size": 12.5, "lines": [
              [seg("1  ", CYAN, 13, bold=True, font=MONO), seg("smithers-orchestrator", WHITE, 13.5, bold=True, font=HEAD), dim("   shipped today · OSS · npm", 11.5)],
              [seg("    the durable JSX workflow runtime · ", DIM, 12, font=BODY), seg("bunx smithers-orchestrator init", GREEN, 12, bold=True, font=MONO)],
              "",
              [seg("2  ", CYAN, 13, bold=True, font=MONO), seg("Smithers (the forge)", WHITE, 13.5, bold=True, font=HEAD), dim("   AGPL · in build · ~78% of MVP", 11.5)],
              [seg("    jj-native code host · landing requests · agent runtime", DIM, 12, font=BODY)],
              [seg("    cloud workspaces on Freestyle VMs · BYOK for LLM keys", DIM, 12, font=BODY)],
              "",
              [seg("3  ", CYAN, 13, bold=True, font=MONO), seg("Smithers GUI", WHITE, 13.5, bold=True, font=HEAD), dim("   native macOS · download today", 11.5)],
              [seg("    embedded Ghostty terminal · time-travel scrubber", DIM, 12, font=BODY)],
              [seg("    picks any agent CLI on PATH (claude / codex / antigravity / kimi / amp / forge)", DIM, 12, font=BODY)],
          ]},
      ])

# 35 — READY TO TRY? (close) --------------------------------------------------
slide(kind="close",
      notes="Try it. B U N X smithers-orchestrator init. Thanks for watching.")


# ════════════════════════════════════════════════════════════════════════════
#  RENDER
# ════════════════════════════════════════════════════════════════════════════
TOTAL = len(SLIDES)


def title_slide(sp, page):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=BG)
    rect(s, 0, 0, SW, 0.06, fill=rgb("0B0F17"))
    txt(s, 0.6, 0.30, 7, 0.3,
        [[seg("smithers", WHITE, 11, bold=True, font=MONO), seg("  ·  technical deck", DIM, 11, font=MONO)]])
    txt(s, SW - 3.1, 0.30, 2.5, 0.3, [[seg(f"{page:02d} / {TOTAL}", DIM, 11, font=MONO)]], align=PP_ALIGN.RIGHT)
    s.shapes.add_picture(os.path.join(ASSETS, "logo.png"), In(5.45), In(1.55), In(1.0), In(1.0))
    txt(s, 0, 2.75, SW, 0.8, [[seg("smithers", WHITE, 46, bold=True, font=MONO)]], align=PP_ALIGN.CENTER)
    txt(s, 0, 3.66, SW, 0.5,
        [[seg("durable AI workflow orchestration as a JSX runtime", CYAN, 17, font=MONO)]], align=PP_ALIGN.CENTER)
    chips = [("9", "primitives", CYAN), ("30+", "components", MAGENTA),
             ("101", "examples", GREEN), ("open", "source", TEAL)]
    cw, gap = 2.0, 0.18
    total_w = len(chips) * cw + (len(chips) - 1) * gap
    x0 = (SW - total_w) / 2
    for i, (num, lab, col) in enumerate(chips):
        x = x0 + i * (cw + gap)
        card(s, x, 4.55, cw, 1.0, fill=PANEL, line=BORDER, radius=0.08)
        txt(s, x, 4.70, cw, 0.45, [[seg(num, col, 21, bold=True, font=MONO)]], align=PP_ALIGN.CENTER)
        txt(s, x, 5.20, cw, 0.3, [[seg(lab, DIM, 11, font=BODY)]], align=PP_ALIGN.CENTER)
    txt(s, 0, 5.95, SW, 0.4,
        [[seg("bunx smithers-orchestrator init", WHITE, 14, bold=True, font=MONO),
          seg("   ·   smithers.sh", DIM, 14, font=MONO)]], align=PP_ALIGN.CENTER)
    notes(s, sp["notes"])


def close_slide(sp, page):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=BG)
    rect(s, 0, 0, SW, 0.06, fill=rgb("0B0F17"))
    txt(s, SW - 3.1, 0.30, 2.5, 0.3, [[seg(f"{page:02d} / {TOTAL}", DIM, 11, font=MONO)]], align=PP_ALIGN.RIGHT)
    s.shapes.add_picture(os.path.join(ASSETS, "logo.png"), In(6.04), In(1.7), In(1.25), In(1.25))
    txt(s, 0, 3.25, SW, 0.7, [[seg("Ready to try?", WHITE, 34, bold=True, font=HEAD)]], align=PP_ALIGN.CENTER)
    card(s, (SW - 6.6) / 2, 4.25, 6.6, 0.95, fill=rgb("13271B"), line=GREEN, radius=0.08)
    txt(s, 0, 4.40, SW, 0.5, [[seg("bunx smithers-orchestrator init", GREEN, 21, bold=True, font=MONO)]], align=PP_ALIGN.CENTER)
    txt(s, 0, 4.92, SW, 0.3, [[seg("scaffolds .smithers/ in any project", DIM, 12, font=BODY)]], align=PP_ALIGN.CENTER)
    txt(s, 0, 5.70, SW, 0.4,
        [[seg("smithers.sh", WHITE, 15, bold=True, font=MONO),
          seg("   ·   docs · llms-full.txt · GUI download", DIM, 15, font=MONO)]], align=PP_ALIGN.CENTER)
    notes(s, sp["notes"])


for i, sp in enumerate(SLIDES):
    page = i + 1
    kind = sp.get("kind", "content")
    if kind == "title":
        title_slide(sp, page)
    elif kind == "close":
        close_slide(sp, page)
    else:
        s = base(sp["tag"], sp["accent"], sp["title"], sp["subtitle"], page, TOTAL)
        render_content(s, sp["blocks"])
        notes(s, sp["notes"])

prs.save(OUT)
print(f"saved {OUT}  ({TOTAL} slides)")
