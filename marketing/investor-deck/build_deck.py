#!/usr/bin/env python3
"""
Build a Google-Slides-compatible .pptx investor deck from the content of
.smithers/workflows/investor.tsx (the terminal pitch deck).

Output: smithers-investor-deck.pptx — import into Google Slides via
"File -> Import slides" or just upload to Drive and "Open with Google Slides".

Every slide's narration from the terminal deck is preserved verbatim as
PowerPoint speaker notes, so the deck is presentation-ready.

Run:  marketing/investor-deck/.venv/bin/python marketing/investor-deck/build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "smithers-investor-deck.pptx")

# ── palette ────────────────────────────────────────────────────────────────
def rgb(h): return RGBColor.from_string(h)
BG      = rgb("0E131C")
PANEL   = rgb("161D2A")
PANEL2  = rgb("1B2433")
BORDER  = rgb("2A3445")
TEXT    = rgb("E6EDF3")
DIM     = rgb("93A1B3")
MUTED   = rgb("64748B")
WHITE   = rgb("FFFFFF")
CYAN    = rgb("56C7E6")
TEAL    = rgb("78D7BA")
GREEN   = rgb("4ADE80")
YELLOW  = rgb("FBBF24")
MAGENTA = rgb("C084FC")
RED     = rgb("F87171")

HEAD = "Inter"
BODY = "Inter"
MONO = "Roboto Mono"

# ── live traction data (mirrors investor.tsx DATA block, as of Jun 2026) ─────
DOWNLOADS = [("Jan", 3235, "3.2k"), ("Feb", 2608, "2.6k"), ("Mar", 2760, "2.8k"),
             ("Apr", 4827, "4.8k"), ("May", 7224, "7.2k"), ("Jun", 16259, "16.3k*")]
STARS = [("Jan", 48, "48"), ("Feb", 76, "76"), ("Mar", 96, "96"),
         ("Apr", 119, "119"), ("May", 232, "232"), ("Jun", 268, "268")]

prs = Presentation()
prs.slide_width = In(13.333)
prs.slide_height = In(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5
_idx = 0
TOTAL = 19

# ── primitives ───────────────────────────────────────────────────────────────
def seg(t, color=TEXT, size=14, bold=False, font=BODY, italic=False):
    return {"t": t, "color": color, "size": size, "bold": bold, "font": font, "italic": italic}

def _no_shadow(sh):
    sh.shadow.inherit = False

def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = slide.shapes.add_shape(shape, In(l), In(t), In(w), In(h))
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sh.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    _no_shadow(sh)
    return sh

def card(slide, l, t, w, h, fill=PANEL, line=BORDER, radius=0.045, line_w=1.0):
    return rect(slide, l, t, w, h, fill=fill, line=line, line_w=line_w,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)

def para(tf, segs, align=PP_ALIGN.LEFT, space_after=4, space_before=0, line_spacing=1.06, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_after is not None: p.space_after = Pt(space_after)
    if space_before: p.space_before = Pt(space_before)
    if line_spacing: p.line_spacing = line_spacing
    for s in segs:
        r = p.add_run(); r.text = s["t"]
        f = r.font
        f.size = Pt(s["size"]); f.bold = s["bold"]; f.italic = s["italic"]
        f.name = s["font"]; f.color.rgb = s["color"]
    return p

def txt(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.06, wrap=True):
    """lines: list of paragraphs, each a list of segs."""
    tb = slide.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, segs in enumerate(lines):
        para(tf, segs, align=align, space_after=space_after, line_spacing=line_spacing, first=(i == 0))
    return tb

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def arrow_glyph(slide, l, t, w, color=MUTED, size=18):
    txt(slide, l, t, w, 0.4, [[seg("→", color, size, bold=True, font=MONO)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# image native size (no Pillow dependency — read PNG/GIF headers) ─────────────
def img_size(path):
    with open(path, "rb") as f:
        head = f.read(32)
    if head[:8] == b"\x89PNG\r\n\x1a\n":
        return int.from_bytes(head[16:20], "big"), int.from_bytes(head[20:24], "big")
    if head[:6] in (b"GIF87a", b"GIF89a"):
        return int.from_bytes(head[6:8], "little"), int.from_bytes(head[8:10], "little")
    return 16, 9

def place_image(slide, name, l, t, w, h, frame=True, caption=None):
    path = os.path.join(ASSETS, name)
    iw, ih = img_size(path)
    ar = iw / ih
    inner = 0.12 if frame else 0.0
    bw, bh = w - 2 * inner, h - 2 * inner
    if caption: bh -= 0.32
    if ar > bw / bh:
        dw = bw; dh = bw / ar
    else:
        dh = bh; dw = bh * ar
    x = l + inner + (bw - dw) / 2
    y = t + inner + (bh - dh) / 2
    if frame:
        card(slide, l, t, w, h, fill=PANEL2, line=BORDER, radius=0.03)
    pic = slide.shapes.add_picture(path, In(x), In(y), In(dw), In(dh))
    pic.line.color.rgb = BORDER; pic.line.width = Pt(0.75)
    if caption:
        txt(slide, l, t + h - 0.30, w, 0.26, [[seg(caption, DIM, 9.5, font=MONO)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return pic

def hbars(slide, l, t, w, h, rows, color, label_w=0.55, value_w=1.45, track=True):
    """rows: (label, value, display) or (label, value, display, color)."""
    n = len(rows)
    maxv = max(r[1] for r in rows) or 1
    gap = 0.10
    rowh = (h - gap * (n - 1)) / n
    bz_l = l + label_w + 0.12
    bz_w = w - label_w - 0.12 - value_w
    for i, r in enumerate(rows):
        lab, val, disp = r[0], r[1], r[2]
        c = r[3] if len(r) > 3 else color
        y = t + i * (rowh + gap)
        txt(slide, l, y, label_w, rowh, [[seg(lab, DIM, 11, font=MONO)]], anchor=MSO_ANCHOR.MIDDLE)
        if track:
            card(slide, bz_l, y + rowh * 0.14, bz_w, rowh * 0.72, fill=PANEL2, line=None, radius=0.5)
        bw = max(0.06, bz_w * (val / maxv))
        card(slide, bz_l, y + rowh * 0.14, bw, rowh * 0.72, fill=c, line=None, radius=0.5)
        txt(slide, bz_l + bw + 0.10, y, value_w, rowh, [[seg(disp, WHITE, 11.5, bold=True, font=MONO)]],
            anchor=MSO_ANCHOR.MIDDLE)

# ── slide chrome ──────────────────────────────────────────────────────────────
TAG_COLOR = {
    "PROBLEM": RED, "SOLUTION": GREEN, "PRODUCT": CYAN, "MOAT": MAGENTA,
    "MARKET": CYAN, "TRACTION": GREEN, "FLYWHEEL": TEAL, "BUSINESS MODEL": YELLOW,
    "COMPETITION": GREEN, "TEAM": CYAN, "GO-TO-MARKET": CYAN, "THE ASK": GREEN, "CLOSE": CYAN,
}

def base(tag, title, subtitle=None, content_top=2.32):
    global _idx
    _idx += 1
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=BG)                       # background
    rect(s, 0, 0, SW, 0.07, fill=rgb("0B0F17"))          # thin top frame
    # brand + page
    txt(s, 0.6, 0.30, 6, 0.3,
        [[seg("smithers", WHITE, 11, bold=True, font=MONO), seg("  ·  investor deck", DIM, 11, font=MONO)]])
    txt(s, SW - 3.1, 0.30, 2.5, 0.3, [[seg(f"{_idx:02d} / {TOTAL}", DIM, 11, font=MONO)]], align=PP_ALIGN.RIGHT)
    # eyebrow tag
    if tag:
        accent = TAG_COLOR.get(tag, CYAN)
        rect(s, 0.6, 0.86, 0.16, 0.16, fill=accent)
        txt(s, 0.86, 0.80, 8, 0.3, [[seg(tag, accent, 12, bold=True, font=MONO)]])
    # title + subtitle
    txt(s, 0.58, 1.16, 12.1, 0.85, [[seg(title, WHITE, 30, bold=True, font=HEAD)]], line_spacing=1.0)
    if subtitle:
        txt(s, 0.6, 1.84, 12.1, 0.45, [[seg(subtitle, DIM, 14.5, font=BODY)]])
    rect(s, 0.6, 2.30 if subtitle else 1.80, 1.7, 0.035, fill=TAG_COLOR.get(tag, CYAN))
    return s

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); _idx += 1
rect(s, 0, 0, SW, SH, fill=BG)
rect(s, 0, 0, SW, SH, fill=None)
# right visual: hero
place_image(s, "hero.png", 7.55, 0.0, 5.78, 7.5, frame=False)
rect(s, 7.55, 0, 0.04, 7.5, fill=BORDER)
# left content
rect(s, 0, 0, 7.55, 7.5, fill=BG)
s.shapes._spTree.remove(s.shapes[-1]._element)  # (left bg already BG; drop extra)
# logo
s.shapes.add_picture(os.path.join(ASSETS, "logo.png"), In(0.85), In(0.85), In(0.95), In(0.95))
txt(s, 1.95, 1.02, 5, 0.7, [[seg("smithers", WHITE, 30, bold=True, font=MONO)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.88, 2.42, 6.55, 1.75,
    [[seg("The durable control plane", WHITE, 33, bold=True, font=HEAD)],
     [seg("for humans + agents", WHITE, 33, bold=True, font=HEAD)],
     [seg("doing real work", WHITE, 33, bold=True, font=HEAD)]],
    line_spacing=1.07, space_after=0)
txt(s, 0.9, 4.5, 6.3, 0.5,
    [[seg("Durable orchestration for background AI agents.", DIM, 15, font=BODY)]])
# metric chips
chips = [("36,913+", "downloads", GREEN), ("268", "GitHub stars", YELLOW),
         ("15+", "ext. contributors", CYAN), ("open", "source", TEAL)]
cw, gap = 1.52, 0.16
for i, (num, lab, col) in enumerate(chips):
    x = 0.88 + i * (cw + gap)
    card(s, x, 5.25, cw, 1.0, fill=PANEL, line=BORDER, radius=0.08)
    txt(s, x, 5.40, cw, 0.45, [[seg(num, col, 19, bold=True, font=MONO)]], align=PP_ALIGN.CENTER)
    txt(s, x, 5.88, cw, 0.3, [[seg(lab, DIM, 10, font=BODY)]], align=PP_ALIGN.CENTER)
txt(s, 0.9, 6.55, 6.3, 0.4, [[seg("smithers.sh", WHITE, 12, bold=True, font=MONO),
                              seg("   ·   github.com/smithersai/smithers", DIM, 12, font=MONO)]])
notes(s, "Hi. This is Smithers. Over the next fifteen minutes I'll walk you through the problem we "
         "solve, the product, the market we're going after, our traction so far, and how we make "
         "money. Short version: every company now wants AI to do real work in the background, and "
         "the plumbing to make that reliable doesn't exist yet. We're building it.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = base("PROBLEM", "Everyone wants AI to do the work.",
         "Almost no one can make it reliable. The gap between a demo and a dependable process.")
rows = [("runs", "once, you're watching", "1000s/week, unattended"),
        ("a crash", "refresh the page", "a dropped customer"),
        ("waits", "seconds", "days — a human must sign off"),
        ("when wrong", "shrug, retry", "money + trust on the line"),
        ("audit", "nobody asks", "who approved this? when?")]
# header row
top = 2.7
txt(s, 0.6, top, 2.0, 0.4, [[seg("", DIM, 12)]])
card(s, 2.7, top, 4.9, 0.55, fill=PANEL, line=GREEN, radius=0.12, line_w=1.25)
txt(s, 2.7, top, 4.9, 0.55, [[seg("AI DEMO", GREEN, 14, bold=True, font=MONO)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
card(s, 7.8, top, 4.9, 0.55, fill=PANEL, line=RED, radius=0.12, line_w=1.25)
txt(s, 7.8, top, 4.9, 0.55, [[seg("REAL BUSINESS PROCESS", RED, 14, bold=True, font=MONO)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
ry = top + 0.72
for lab, a, b in rows:
    card(s, 0.6, ry, 2.0, 0.6, fill=PANEL2, line=None, radius=0.12)
    txt(s, 0.75, ry, 1.85, 0.6, [[seg(lab, DIM, 12.5, font=MONO)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 2.85, ry, 4.6, 0.6, [[seg(a, TEXT, 13, font=BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 7.95, ry, 4.6, 0.6, [[seg(b, rgb("FCA5A5"), 13, font=BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    ry += 0.68
card(s, 0.6, ry + 0.05, 12.1, 0.62, fill=rgb("1F1518"), line=YELLOW, radius=0.08)
txt(s, 0.6, ry + 0.05, 12.1, 0.62,
    [[seg("The AI got good. The plumbing to run it like a business didn't.", YELLOW, 15, bold=True, font=BODY)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "The problem. Every small and medium business has dozens of multi-step processes that are "
         "begging to be automated. And for the first time, AI agents are good enough to actually do "
         "them. But there's a gap. A demo that works once in a chat window is not a process you can "
         "run a thousand times a week, unattended, where a single crash doesn't drop a customer's "
         "refund on the floor. Going from a cool demo to something you'd actually trust with the "
         "business is where everyone gets stuck.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHY THIS IS HARD
# ════════════════════════════════════════════════════════════════════════════
s = base("PROBLEM", "Why this is hard (and stays hard)",
         "Three traps: a moving target, a substrate nobody builds, one-size-fits-all tools.")
traps = [
    (YELLOW, "TRAP 1", "The moving target",
     ["chains → ReAct → tools → plan-execute →", "crews → background agents → ?",
      "", "Couple your business to this year's topology,", "and you rebuild it every 6 months."]),
    (RED, "TRAP 2", "The substrate nobody wants to build",
     ["✗  durable steps + resume", "✗  retries + crash recovery", "✗  pause for a human, free",
      "✗  a real audit trail", "", "A queue + DB gets you ~60%, badly."]),
    (MAGENTA, "TRAP 3", "Tools that exist are one-size-fits-all",
     ["Off-the-shelf orchestrators ship someone", "else's pipeline, baked in.", "",
      "Your process isn't their process —", "and you can't reshape theirs."]),
]
cw = 3.95; gap = 0.12; x0 = 0.6; ty = 2.85
for i, (col, kick, title, lines) in enumerate(traps):
    x = x0 + i * (cw + gap)
    card(s, x, ty, cw, 3.95, fill=PANEL, line=BORDER, radius=0.05)
    rect(s, x, ty, cw, 0.07, fill=col)
    txt(s, x + 0.28, ty + 0.30, cw - 0.5, 0.3, [[seg(kick, col, 12, bold=True, font=MONO)]])
    txt(s, x + 0.28, ty + 0.66, cw - 0.5, 0.65, [[seg(title, WHITE, 16.5, bold=True, font=HEAD)]], line_spacing=1.0)
    body = [[seg(l, RED if l.startswith("✗") else DIM, 12.5, font=(MONO if l.startswith("✗") else BODY))] for l in lines]
    txt(s, x + 0.28, ty + 1.55, cw - 0.5, 2.2, body, space_after=3, line_spacing=1.12)
notes(s, "Why is this hard? Three reasons. First, the right way to build an AI agent changes every "
         "six months. Chains, then tools, then crews, then background agents. If a business couples "
         "its automation to this year's fashion, they rebuild it next year. Second, the reliability "
         "layer underneath — retries, resume-after-crash, pausing for a human, an audit trail — is "
         "about sixty percent of a real orchestration engine, and almost everyone tries to rebuild "
         "it by hand on top of a queue and a database, and gets it wrong. And third, the "
         "orchestration tools that do exist are one-size-fits-all. They ship someone else's opinion "
         "of the pipeline, baked in, and your business process is never quite their business "
         "process. The substrate, fitted to your problem, is exactly what we sell.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE SOLUTION (3-layer stack)
# ════════════════════════════════════════════════════════════════════════════
s = base("SOLUTION", "Be the layer that doesn't change.",
         "Three layers, three speeds of change. We sell the bottom one — the part you never throw away.")
layers = [
    (RED, "MODEL LAYER", "volatile · changes weekly", "GPT · Claude · Gemini · Kimi · whatever wins this month", 0.42, False),
    (YELLOW, "AGENT / TOPOLOGY LAYER", "fluid · changes quarterly", "ReAct · crew · swarm · background agents · the next fad", 0.42, False),
    (GREEN, "ORCHESTRATION LAYER  —  this is Smithers", "stable · you build on it once", "durable steps · retries · state · human-in-the-loop · audit", 1.0, True),
]
ly = 2.95; lh = 1.12; lx = 0.6; lw = 12.13
for col, name, speed, detail, alpha, hot in layers:
    card(s, lx, ly, lw, lh, fill=(rgb("13271B") if hot else PANEL), line=(col if hot else BORDER), radius=0.05, line_w=(1.75 if hot else 1.0))
    rect(s, lx, ly, 0.12, lh, fill=col)
    txt(s, lx + 0.4, ly + 0.18, 8.5, 0.4, [[seg(name, col if hot else TEXT, 17 if hot else 15.5, bold=True, font=HEAD)]])
    txt(s, lx + 0.4, ly + 0.62, 8.5, 0.4, [[seg(detail, DIM, 12.5, font=MONO)]])
    txt(s, lx + lw - 3.4, ly, 3.0, lh, [[seg(speed, col, 12.5, bold=hot, font=BODY)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    ly += lh + 0.18
txt(s, 0.6, ly + 0.06, 12.13, 0.5,
    [[seg("We sell the bottom layer — the one a business can't live without, and least wants to maintain itself.", DIM, 13.5, font=BODY)]],
    align=PP_ALIGN.CENTER)
notes(s, "Our solution is to be the layer that doesn't change. Underneath every trendy agent pattern "
         "there's a stable foundation: steps, state, retries, waiting, and an audit trail. Smithers "
         "is that foundation, shipped as a tool a developer can drop into any project in one command. "
         "The fashions on top can change every quarter. We're the part you build on and never throw away.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SO, WHAT IS SMITHERS? (node graph)
# ════════════════════════════════════════════════════════════════════════════
s = base("PRODUCT", "So, what is Smithers?",
         "Wire humans and AI agents together as steps in one long-running, crash-proof workflow.")
# node graph
nodes = [("agent", CYAN), ("agent", CYAN), ("human\napproval", YELLOW), ("agent", CYAN), ("done", GREEN)]
nx = 0.85; ny = 2.95; nw = 1.95; nh = 1.0; ngap = 0.42
for i, (lab, col) in enumerate(nodes):
    x = nx + i * (nw + ngap)
    card(s, x, ny, nw, nh, fill=PANEL, line=col, radius=0.12, line_w=1.5)
    lines = [[seg(part, col, 14.5, bold=True, font=MONO)] for part in lab.split("\n")]
    txt(s, x, ny, nw, nh, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
    if i < len(nodes) - 1:
        arrow_glyph(s, x + nw, ny, ngap, color=MUTED, size=20)
txt(s, 0.6, ny + nh + 0.22, 12.13, 0.35,
    [[seg("intake → triage → draft → human approval → execute → notify     —  each box is a durable node", DIM, 12.5, font=MONO)]],
    align=PP_ALIGN.CENTER)
# guarantees + screenshot
gy = 4.85
gtxt = [
    [seg("✓  ", GREEN, 14, bold=True), seg("runs for hours or days", TEXT, 13.5)],
    [seg("✓  ", GREEN, 14, bold=True), seg("survives a crash mid-step", TEXT, 13.5)],
    [seg("✓  ", GREEN, 14, bold=True), seg("pauses for a human, for free", TEXT, 13.5)],
    [seg("✓  ", GREEN, 14, bold=True), seg("every step is auditable", TEXT, 13.5)],
]
card(s, 0.6, gy, 6.05, 2.05, fill=PANEL, line=BORDER, radius=0.05)
txt(s, 0.92, gy + 0.26, 5.5, 1.6, gtxt, space_after=9, line_spacing=1.0)
txt(s, 0.92, gy + 1.62, 5.6, 0.4,
    [[seg("Humans and agents are the same primitive: a step that can fail and resume.", TEAL, 11.5, italic=True, font=BODY)]])
place_image(s, "ask-human.png", 6.95, gy, 5.78, 2.05, frame=True, caption="human-in-the-loop approval — a paused workflow is a row, not a server")
notes(s, "So concretely, what is Smithers? Smithers lets a company wire humans and AI agents together "
         "as steps in a single, long-running workflow that survives crashes, deploys, and overnight "
         "waits for a human approval. Think of it as the durable control plane underneath all the AI "
         "automation a business is about to build.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOW IT WORKS (the loop)
# ════════════════════════════════════════════════════════════════════════════
s = base("PRODUCT", "How it works",
         "A workflow is a tree of steps. State lives in a database. Crash → re-read → resume.")
loop = [("read\ntree", CYAN), ("run the\nready steps", CYAN), ("write\nresults", CYAN), ("re-read\nwith state", CYAN)]
lx = 0.85; ly = 2.85; lw = 2.5; lh = 0.95; lg = 0.55
for i, (lab, col) in enumerate(loop):
    x = lx + i * (lw + lg)
    card(s, x, ly, lw, lh, fill=PANEL, line=col, radius=0.1, line_w=1.4)
    txt(s, x, ly, lw, lh, [[seg(p, col, 13.5, bold=True, font=MONO)] for p in lab.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
    if i < 3:
        arrow_glyph(s, x + lw, ly, lg, color=MUTED, size=18)
# return loop line
rect(s, lx + 0.4, ly + lh + 0.42, lw * 4 + lg * 3 - 0.8, 0.025, fill=MUTED)
rect(s, lx + 0.4, ly + lh + 0.10, 0.025, 0.34, fill=MUTED)
rect(s, lx + lw * 4 + lg * 3 - 0.4, ly + lh + 0.10, 0.025, 0.34, fill=MUTED)
txt(s, lx, ly + lh + 0.50, lw * 4 + lg * 3, 0.3,
    [[seg("state in a database, not in memory  →  a crash just re-reads and continues", DIM, 12, font=MONO)]],
    align=PP_ALIGN.CENTER)
# benefits + proof image
by = 4.55
bene = [("Crash recovery", "kill it anywhere, resume from the last saved step"),
        ("Pause for humans", "a waiting workflow is a row, not a server — $0"),
        ("Audit trail", "every step, attempt, and approval is queryable"),
        ("Time travel", "fork any past state, replay, diff — git for work")]
card(s, 0.6, by, 6.05, 2.4, fill=PANEL, line=BORDER, radius=0.05)
yy = by + 0.22
for t1, t2 in bene:
    txt(s, 0.92, yy, 5.5, 0.55,
        [[seg("✓ ", GREEN, 12.5, bold=True), seg(t1 + "   ", WHITE, 13, bold=True, font=HEAD), seg(t2, DIM, 11.5, font=BODY)]],
        line_spacing=1.0)
    yy += 0.56
place_image(s, "runs-live.gif", 6.95, by, 5.78, 2.4, frame=True, caption="a run executing live — state in the DB, resumes after a crash (animated)")
notes(s, "Here's how it works, mechanically, without the jargon. You describe a workflow as a tree of "
         "steps. Some steps are agents, some are humans, some are plain code. Smithers runs the steps "
         "that are ready, writes each result to a database, and re-reads the tree against the new "
         "state. That loop is the whole engine. Because the state lives in the database and not in "
         "memory, a crash isn't a catastrophe — you just re-read the database and keep going from the "
         "last finished step. And because every step is a row, you get a complete audit trail and "
         "even time-travel for free.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — THE PRODUCT (compose / run)
# ════════════════════════════════════════════════════════════════════════════
s = base("PRODUCT", "The product — compose, then run like production",
         "A toolbox of proven blocks, plus the unsexy 60% nobody wants to build.")
compose = [("Nodes", "agent · human · code — one durable step"),
           ("Control flow", "sequence · parallel · branch · loop"),
           ("Human-in-loop", "approval gates · ask-a-human · wait"),
           ("Pattern catalog", "review loops · panels · debate · saga"),
           ("Quality + memory", "scorers + evals · facts across runs"),
           ("Provider-neutral", "Claude · GPT · Gemini · Kimi — swap")]
run = [("Operate", "Docker-Compose-style CLI · auto-resume supervisor"),
       ("Observe", "metrics · traces · one-command Grafana stack"),
       ("Schedule", "durable cron + alerts that survive restarts"),
       ("Cost control", "per-workflow token / $ budgets — warn or stop"),
       ("Time travel", "fork · replay · diff · hot-edit live runs"),
       ("Serve", "ships its own API + UI — what 3rd parties build on")]
def feature_col(x, w, title, col, items):
    card(s, x, 2.8, w, 3.55, fill=PANEL, line=BORDER, radius=0.04)
    rect(s, x, 2.8, w, 0.07, fill=col)
    txt(s, x + 0.3, 3.02, w - 0.6, 0.4, [[seg(title, col, 14.5, bold=True, font=HEAD)]])
    yy = 3.52
    for k, v in items:
        txt(s, x + 0.3, yy, w - 0.6, 0.5,
            [[seg(k, WHITE, 12.5, bold=True, font=MONO)]], space_after=1)
        txt(s, x + 0.3, yy + 0.22, w - 0.6, 0.3, [[seg(v, DIM, 11, font=BODY)]])
        yy += 0.47
feature_col(0.6, 6.0, "COMPOSE — proven building blocks", CYAN, compose)
feature_col(6.73, 6.0, "RUN IN PRODUCTION — the 60% nobody builds", GREEN, run)
card(s, 0.6, 6.55, 12.13, 0.62, fill=rgb("13271B"), line=GREEN, radius=0.08)
txt(s, 0.6, 6.55, 12.13, 0.62,
    [[seg("We rebuilt a popular agent tool on these blocks → ", TEXT, 14, font=BODY),
      seg("~80% less code.", GREEN, 14, bold=True, font=BODY)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Let me go a little deeper on the product, because the depth is the moat, in two halves. On "
         "the build side, Smithers is a toolbox, not a rigid framework. Three kinds of node — an "
         "agent, a human, or plain code, all the same durable step. Control flow, human-in-the-loop "
         "gates, and a catalog of named patterns shipped as drop-in components: review loops, panels, "
         "debates, supervisors, sagas. On the run side, it's real production infrastructure: a "
         "Docker-Compose-style command line, one-command observability, durable scheduling and "
         "alerts, per-workflow cost budgets, isolation, time travel, and the ability for a workflow "
         "to serve its own API and web interface. When we rebuilt a popular agent tool on these "
         "blocks, we cut about eighty percent of its code.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHY AGENTS BUILD BETTER (React) + benchmarks
# ════════════════════════════════════════════════════════════════════════════
s = base("MOAT", "Why agents build better on Smithers",
         "We mapped orchestration onto React/JSX — the one paradigm every model already knows.")
# insight card
card(s, 0.6, 2.8, 5.55, 3.95, fill=PANEL, line=MAGENTA, radius=0.05)
rect(s, 0.6, 2.8, 5.55, 0.07, fill=MAGENTA)
txt(s, 0.9, 3.05, 5.0, 0.4, [[seg("THE INSIGHT", MAGENTA, 12.5, bold=True, font=MONO)]])
txt(s, 0.9, 3.45, 5.0, 3.1, [
    [seg("Agents — and humans — are saturated with React / JSX: the most heavily represented way to build software in training data.", TEXT, 14, font=BODY)],
    [seg("", TEXT, 6)],
    [seg("Orchestration is a tree-of-steps problem. React composes trees. So we mapped durable orchestration directly onto JSX.", DIM, 13, font=BODY)],
    [seg("", TEXT, 6)],
    [seg("Payoff: agents write the paradigm they know best — so they produce more correct, more complex graphs than imperative code or a DSL.", GREEN, 13, bold=True, font=BODY)],
], space_after=4, line_spacing=1.14)
# benchmarks card
card(s, 6.35, 2.8, 6.38, 3.95, fill=PANEL, line=BORDER, radius=0.05)
txt(s, 6.65, 3.0, 5.8, 0.4, [[seg("BENCHMARKS", CYAN, 12.5, bold=True, font=MONO),
                              seg("   full results release next week", DIM, 11, font=BODY)]])
txt(s, 6.65, 3.45, 5.8, 0.3, [[seg("SWE-EVO — long-horizon (build a release from its notes)", DIM, 11.5, font=BODY)]])
hbars(s, 6.65, 3.72, 5.8, 0.78,
      [("frontier", 23, "~23%", RED), ("Smithers", 71, "71%  (dvc)", GREEN)], GREEN, label_w=0.9, value_w=1.55)
txt(s, 6.65, 4.66, 5.8, 0.3, [[seg("RoadmapBench — completion (multi-target upgrades)", DIM, 11.5, font=BODY)]])
hbars(s, 6.65, 4.93, 5.8, 0.78,
      [("SOTA", 69, "0.69", RED), ("Smithers", 86, "0.86", GREEN)], GREEN, label_w=0.9, value_w=1.55)
txt(s, 6.65, 5.92, 5.8, 0.7, [
    [seg("+ Claw-Eval-Live", WHITE, 11.5, bold=True, font=MONO), seg("  — 105 real enterprise workflows, lands next week", DIM, 11, font=BODY)],
    [seg("+ SWE-Bench Pro", WHITE, 11.5, bold=True, font=MONO), seg("  — 731 professional tasks, same harness", DIM, 11, font=BODY)],
], space_after=3)
notes(s, "Here's a non-obvious insight that's becoming one of our biggest advantages. Agents, and "
         "honestly humans too, are saturated with React and JSX. It's the most heavily represented "
         "way to build software in any model's training data. Orchestration is fundamentally a "
         "tree-of-steps problem, and React is a language for composing trees. So we mapped the entire "
         "orchestration problem onto JSX. The payoff is that agents produce far more correct and more "
         "complex orchestration graphs on Smithers than on imperative code or a DSL. Early numbers: "
         "on SWE-EVO, where frontier agents resolve only about twenty to twenty-five percent, "
         "Smithers resolves seventy-one percent of the dvc subset. On RoadmapBench we're already "
         "above the public state of the art. And Claw-Eval-Live, a hundred and five real enterprise "
         "workflows, lands next week.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WHY IT'S STICKY
# ════════════════════════════════════════════════════════════════════════════
s = base("MOAT", "Why it's sticky",
         "We don't automate one task. We become the fabric the work flows through.")
card(s, 0.6, 2.95, 6.0, 3.0, fill=PANEL, line=RED, radius=0.05)
rect(s, 0.6, 2.95, 6.0, 0.07, fill=RED)
txt(s, 0.9, 3.2, 5.4, 0.4, [[seg("A POINT AI TOOL", RED, 13, bold=True, font=MONO)]])
for i, l in enumerate(["automates one task", "one team, one box", "easy to swap out"]):
    txt(s, 0.9, 3.75 + i * 0.62, 5.4, 0.5, [[seg("•  ", RED, 13), seg(l, TEXT, 14, font=BODY)]])
card(s, 6.73, 2.95, 6.0, 3.0, fill=rgb("13271B"), line=GREEN, radius=0.05)
rect(s, 6.73, 2.95, 6.0, 0.07, fill=GREEN)
txt(s, 7.03, 3.2, 5.4, 0.4, [[seg("SMITHERS", GREEN, 13, bold=True, font=MONO)]])
for i, l in enumerate(["orchestrates whole processes", "humans + agents across teams", "the rails everything else runs on"]):
    txt(s, 7.03, 3.75 + i * 0.62, 5.4, 0.5, [[seg("•  ", GREEN, 13), seg(l, WHITE, 14, bold=True, font=BODY)]])
card(s, 0.6, 6.2, 12.13, 0.95, fill=PANEL, line=YELLOW, radius=0.06)
txt(s, 0.6, 6.28, 12.13, 0.4, [[seg("Land one workflow. Expand to how the company runs.", YELLOW, 16, bold=True, font=HEAD)]], align=PP_ALIGN.CENTER)
txt(s, 0.6, 6.72, 12.13, 0.35, [[seg("As more processes move onto Smithers, switching cost compounds. Ripping us out means re-plumbing the company.", DIM, 12, font=BODY)]], align=PP_ALIGN.CENTER)
notes(s, "Here's why we think this gets sticky. Most AI tools automate one task. We orchestrate "
         "cross-functional collaboration — humans and agents as nodes in larger, durable workflows "
         "that span teams. Once a company's refund process, its onboarding, its incident response, "
         "its content pipeline all run as Smithers workflows, we're not a feature they can swap out. "
         "We're the fabric the work flows through. Ripping us out means re-plumbing how the company "
         "operates.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THE MARKET
# ════════════════════════════════════════════════════════════════════════════
s = base("MARKET", "The market",
         "Small + medium businesses automating their workflows — the most process work, the least platform team.")
cards3 = [(CYAN, "WHO", "SMBs with lots of multi-step process work and no platform team to build reliable automation in-house."),
          (YELLOW, "PAIN", "Every process is manual, or a brittle script that breaks silently. They can't hire a Temporal team."),
          (GREEN, "NOW", "AI agents are finally good enough to DO the steps. Reliability is the missing piece — and that's what we sell.")]
cw = 3.95; x0 = 0.6
for i, (col, k, v) in enumerate(cards3):
    x = x0 + i * (cw + 0.12)
    card(s, x, 2.85, cw, 2.0, fill=PANEL, line=BORDER, radius=0.05)
    rect(s, x, 2.85, cw, 0.07, fill=col)
    txt(s, x + 0.28, 3.08, cw - 0.5, 0.4, [[seg(k, col, 14, bold=True, font=MONO)]])
    txt(s, x + 0.28, 3.52, cw - 0.5, 1.3, [[seg(v, TEXT, 13, font=BODY)]], line_spacing=1.16)
card(s, 0.6, 5.05, 12.13, 1.95, fill=rgb("11202B"), line=CYAN, radius=0.05)
txt(s, 0.9, 5.28, 11.5, 0.4, [[seg("The unlock:  an SMB doesn't need to hire engineers to adopt us.", CYAN, 16, bold=True, font=HEAD)]])
txt(s, 0.9, 5.74, 11.5, 0.4, [[seg("An agent reads our docs and writes the workflow. We designed for exactly that.", DIM, 13.5, font=BODY)]])
txt(s, 0.9, 6.22, 11.5, 0.7, [
    [seg("Demand is already crossing the technical fence — ", YELLOW, 13.5, bold=True, font=BODY),
     seg("non-technical users ask us for a Zapier / n8n replacement on Smithers.", TEXT, 13.5, font=BODY)],
    [seg("That's the SMB workflow-automation market asking for us by name.", DIM, 12.5, italic=True, font=BODY)],
], space_after=3)
notes(s, "The market we're going after is small and medium businesses automating their workflows and "
         "making them more efficient. This is the part of the economy that has the most repetitive, "
         "multi-step process work and the least ability to staff a platform team to build reliable "
         "automation in-house. They can't hire a Temporal team. But they can adopt a tool, and "
         "increasingly they can point an AI agent at that tool and have it build the workflow for "
         "them. That's our entry point.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TRACTION (charts)
# ════════════════════════════════════════════════════════════════════════════
s = base("TRACTION", "Traction — all real numbers, pulled this week",
         "36,913 downloads · 268 stars · 15+ external contributors · launched Jan 2026.")
# downloads chart
card(s, 0.6, 2.78, 6.0, 3.45, fill=PANEL, line=BORDER, radius=0.04)
txt(s, 0.9, 3.0, 5.4, 0.35, [[seg("npm downloads · monthly", GREEN, 14, bold=True, font=HEAD)]])
hbars(s, 0.9, 3.5, 5.4, 2.0, DOWNLOADS, GREEN, label_w=0.5, value_w=1.55)
txt(s, 0.9, 5.62, 5.4, 0.5, [[seg("all-time 36,913 · last 30d 21,001", WHITE, 11.5, bold=True, font=MONO)],
                             [seg("* June month-to-date — 2.2× May, still climbing", DIM, 11, font=BODY)]], space_after=2)
# stars chart
card(s, 6.73, 2.78, 6.0, 3.45, fill=PANEL, line=BORDER, radius=0.04)
txt(s, 7.03, 3.0, 5.4, 0.35, [[seg("GitHub stars · cumulative", YELLOW, 14, bold=True, font=HEAD)]])
hbars(s, 7.03, 3.5, 5.4, 2.0, STARS, YELLOW, label_w=0.5, value_w=1.2)
txt(s, 7.03, 5.62, 5.4, 0.5, [[seg("268 stars · 30 forks · 15+ external contributors", WHITE, 11.5, bold=True, font=MONO)],
                              [seg("5.6× since launch · open source = distribution engine", DIM, 11, font=BODY)]], space_after=2)
txt(s, 0.6, 6.42, 12.13, 0.3, [[seg("source: npmjs.org + github.com/smithersai/smithers · as of Jun 2026", MUTED, 10.5, font=MONO)]], align=PP_ALIGN.CENTER)
notes(s, "Now, traction — all real numbers, pulled this week. We launched on NPM in January. "
         "Downloads: about thirty-seven thousand all-time, twenty-one thousand in just the last "
         "thirty days, and June is already more than double May with the month not over. That's the "
         "hockey stick you want — accelerating, not flattening. The community tracks it: two hundred "
         "sixty-eight GitHub stars, up more than five-x since launch, thirty forks, and fifteen-plus "
         "external contributors across three thousand commits. Open source is our distribution "
         "engine: every star is a developer who already trusts the tool.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — WHAT USERS ARE TELLING US
# ════════════════════════════════════════════════════════════════════════════
s = base("TRACTION", "What users are telling us",
         "A step-function moment — and demand from people who can't code.")
q1 = [("“A Claude-3.5 step-function moment for my own productivity.”", GREEN, "— how multiple users describe adopting Smithers")]
q2 = [("“Can you build me a Zapier / n8n replacement on this?”", CYAN, "— non-technical users, unprompted")]
qy = 2.85
for quote, col, attr in [q1[0], q2[0]]:
    card(s, 0.6, qy, 12.13, 1.35, fill=PANEL, line=BORDER, radius=0.05)
    rect(s, 0.6, qy, 0.1, 1.35, fill=col)
    txt(s, 1.0, qy + 0.22, 11.4, 0.6, [[seg(quote, WHITE, 19, bold=True, font=HEAD)]], line_spacing=1.0)
    txt(s, 1.0, qy + 0.92, 11.4, 0.3, [[seg(attr, DIM, 12.5, italic=True, font=BODY)]])
    qy += 1.55
card(s, 0.6, qy + 0.0, 12.13, 1.5, fill=rgb("13271B"), line=GREEN, radius=0.05)
txt(s, 0.95, qy + 0.18, 11.5, 0.4, [[seg("Why it matters", GREEN, 13.5, bold=True, font=MONO)]])
for i, l in enumerate([
    "The wedge from developers → operators → whole orgs is happening organically.",
    "Pull from non-technical users = the SMB market reaching past the dev seat.",
    "“Step-function moment” is the language of retention, not curiosity."]):
    txt(s, 0.95, qy + 0.55 + i * 0.30, 11.5, 0.3, [[seg("▸ ", GREEN, 12), seg(l, TEXT, 12.5, font=BODY)]])
notes(s, "The qualitative signal matches the numbers, and it's the part that gives me the most "
         "conviction. Users describe Smithers as a step-function moment in their own productivity. "
         "One called it their Claude three-point-five moment. And the demand is jumping the technical "
         "fence. Non-technical people are asking us to build them a Zapier or an n8n replacement on "
         "top of Smithers. That's the exact small-and-medium-business automation market, asking for "
         "us by name, before we've spent a dollar on sales.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — OPEN-SOURCE FLYWHEEL
# ════════════════════════════════════════════════════════════════════════════
s = base("FLYWHEEL", "The open-source flywheel",
         "Free adoption → contribution + signal → better product — and others build on us.")
fly = [("adoption", GREEN), ("contribute\n+ issues", CYAN), ("roadmap\nsignal", YELLOW), ("better\nproduct", MAGENTA)]
fx = 0.85; fy = 2.85; fw = 2.45; fh = 0.95; fg = 0.5
for i, (lab, col) in enumerate(fly):
    x = fx + i * (fw + fg)
    card(s, x, fy, fw, fh, fill=PANEL, line=col, radius=0.1, line_w=1.4)
    txt(s, x, fy, fw, fh, [[seg(p, col, 13.5, bold=True, font=MONO)] for p in lab.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
    if i < 3:
        arrow_glyph(s, x + fw, fy, fg, color=MUTED, size=18)
txt(s, fx, fy + fh + 0.12, fw * 4 + fg * 3, 0.3,
    [[seg("…which drives more adoption (21,001/mo). Distribution cost trends to zero.", DIM, 12, font=MONO)]],
    align=PP_ALIGN.CENTER)
# proof + gif
py = 4.5
card(s, 0.6, py, 6.4, 2.55, fill=PANEL, line=BORDER, radius=0.05)
txt(s, 0.9, py + 0.2, 5.8, 0.4, [[seg("Proof it's working — built on Smithers, by people who don't work here", CYAN, 13, bold=True, font=HEAD)]], line_spacing=1.0)
for i, (k, v) in enumerate([("Third-party UIs", "dashboards + front-ends on the engine"),
                            ("Kubernetes automation", "real workflows orchestrated in prod clusters"),
                            ("Custom workflow packs", "domain patterns we never wrote")]):
    txt(s, 0.9, py + 0.75 + i * 0.42, 5.9, 0.4, [[seg("▸ ", CYAN, 12.5, bold=True), seg(k + "  ", WHITE, 13, bold=True, font=MONO), seg(v, DIM, 11, font=BODY)]])
txt(s, 0.9, py + 2.06, 5.9, 0.4, [[seg("Third parties run our R&D in production, free — and show us what enterprises will pay for.", TEAL, 11.5, italic=True, font=BODY)]], line_spacing=1.05)
place_image(s, "chat-first.gif", 7.3, py, 5.43, 2.55, frame=True, caption="live: humans + agents working side by side (animated)")
notes(s, "Open source isn't just the vibe, it's the business engine, and it's a flywheel. Free "
         "adoption brings developers in. Some contribute code and file issues; all of them generate "
         "signal about what real workloads need, which directs the roadmap better than any focus "
         "group. The clearest proof it's working: third parties are already building on us. People we "
         "don't employ have built their own user interfaces on our engine, and teams are automating "
         "workflows inside Kubernetes. Other people are running our R&D in production, for free, and "
         "showing us exactly where to go.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — HOW WE MAKE MONEY
# ════════════════════════════════════════════════════════════════════════════
s = base("BUSINESS MODEL", "How we make money",
         "At-cost collaboration tool drives daily usage; enterprise features drive revenue.")
tiers = [(GREEN, "1", "OPEN SOURCE ENGINE", "free · drives adoption + the flywheel", "the durable orchestrator on npm"),
         (CYAN, "2", "COLLABORATION TOOL", "at cost · we control providers, no token markup", "Amp-Code-like · cloud workspaces + native app · cheapest credible default"),
         (YELLOW, "3", "ENTERPRISE — the revenue line", "BYOK · SSO · governance · audit · on-prem · support / SLAs", "money is in governance + control, not tokens")]
ty = 2.85
for col, n, name, mid, sub in tiers:
    h = 1.18 if n != "3" else 1.18
    hot = n == "3"
    card(s, 0.6, ty, 8.4, h, fill=(rgb("1F1B10") if hot else PANEL), line=(col if hot else BORDER), radius=0.05, line_w=(1.6 if hot else 1.0))
    rect(s, 0.6, ty, 0.12, h, fill=col)
    txt(s, 0.95, ty, 0.7, h, [[seg(n, col, 30, bold=True, font=MONO)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.7, ty + 0.18, 7.0, 0.4, [[seg(name, col if hot else WHITE, 16, bold=True, font=HEAD)]])
    txt(s, 1.7, ty + 0.58, 7.0, 0.3, [[seg(mid, DIM, 12, font=MONO)]])
    txt(s, 1.7, ty + 0.86, 7.0, 0.3, [[seg(sub, MUTED, 11, font=BODY)]])
    ty += h + 0.13
# funnel side note + usage image
place_image(s, "usage.png", 9.2, 2.85, 3.53, 3.69, frame=True, caption="per-account usage + cost control")
card(s, 0.6, ty + 0.02, 8.4, 0.66, fill=PANEL, line=None, radius=0.06)
txt(s, 0.6, ty + 0.02, 8.4, 0.66, [[seg("Adoption is free. Usage is at-cost. Money is in governance + control.", YELLOW, 13.5, bold=True, font=BODY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "How we make money. We're still discovering the exact model, but here's the current shape. "
         "We're building a collaboration tool, similar in spirit to Amp Code, where we control the AI "
         "providers and offer the AI compute at cost — no markup on tokens. That makes us the obvious "
         "default for individuals and small teams. Then we monetize the enterprise: "
         "bring-your-own-key, single sign-on, governance, audit, on-premise, support. The open-source "
         "engine drives adoption. The at-cost tool drives daily usage. The enterprise features drive "
         "revenue.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — COMPETITIVE LANDSCAPE
# ════════════════════════════════════════════════════════════════════════════
s = base("COMPETITION", "The competitive landscape",
         "Everyone looks like a competitor if you squint. None sit where we sit.")
comp = [("Temporal", "durable-workflow leader", "built for platform engineers writing code — not AI-native, no human/agent nodes"),
        ("Model-lab workflow tools", "e.g. Claude workflows", "one topology, one provider — lock-in. We're provider-neutral"),
        ("Amp Code + coding agents", "great products", "but a product, not a substrate you compose on")]
cy = 2.82
for name, tag, desc in comp:
    card(s, 0.6, cy, 12.13, 1.0, fill=PANEL, line=BORDER, radius=0.05)
    txt(s, 0.9, cy + 0.16, 4.5, 0.7, [[seg(name, WHITE, 15, bold=True, font=HEAD)],
                                       [seg(tag, MUTED, 11, italic=True, font=BODY)]], space_after=2)
    txt(s, 5.3, cy, 7.2, 1.0, [[seg(desc, DIM, 13, font=BODY)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    cy += 1.12
card(s, 0.6, cy + 0.0, 12.13, 1.05, fill=rgb("13271B"), line=GREEN, radius=0.05)
txt(s, 0.9, cy + 0.16, 11.5, 0.4, [[seg("Smithers", GREEN, 16, bold=True, font=HEAD),
                                    seg("  = open · durable · provider-neutral · humans + agents as nodes · agent-authorable", WHITE, 14, font=BODY)]])
txt(s, 0.9, cy + 0.6, 11.5, 0.4, [[seg("Differentiation is case-by-case — happy to go competitor by competitor.", DIM, 12.5, italic=True, font=BODY)]])
notes(s, "On competition. If you squint, everyone looks like a competitor. Temporal is the "
         "durable-workflow leader, but it's built for platform engineers writing code, not for "
         "AI-native, human-plus-agent workflows. The model labs ship their own workflow tools, but "
         "those are bets on one topology and they lock you to one provider. Amp Code and the coding "
         "agents overlap on the collaboration surface, but they're products, not a substrate you "
         "build on. Our position is the open, durable substrate where humans and agents are "
         "first-class nodes, provider-neutral, and agent-authorable.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — WHY THIS FOUNDER
# ════════════════════════════════════════════════════════════════════════════
s = base("TEAM", "Why this founder",
         "A proven agentic coder building the tool that makes agentic coding a commodity.")
card(s, 0.6, 2.85, 6.0, 3.95, fill=PANEL, line=BORDER, radius=0.05)
rect(s, 0.6, 2.85, 6.0, 0.07, fill=CYAN)
txt(s, 0.9, 3.08, 5.4, 0.4, [[seg("TRACK RECORD", CYAN, 12.5, bold=True, font=MONO)]])
tr = [
    ("Proven operator", "Shipped hard production code on far weaker models than today's."),
    ("Fastest EVM ever", "Beat revm — the previous leader — built with Claude 3.5."),
    ("Ex-Google, applied", "Smithers' high-throughput agent organization borrows Google's monorepo discipline for coordinating work at scale."),
]
yy = 3.5
for k, v in tr:
    txt(s, 0.9, yy, 5.4, 0.35, [[seg(k, WHITE, 13.5, bold=True, font=HEAD)]])
    txt(s, 0.9, yy + 0.30, 5.4, 0.7, [[seg(v, DIM, 12, font=BODY)]], line_spacing=1.12)
    yy += 1.08
card(s, 6.73, 2.85, 6.0, 3.95, fill=rgb("1F1B10"), line=YELLOW, radius=0.05)
rect(s, 6.73, 2.85, 6.0, 0.07, fill=YELLOW)
txt(s, 7.03, 3.08, 5.4, 0.4, [[seg("THE THESIS", YELLOW, 12.5, bold=True, font=MONO)]])
txt(s, 7.03, 3.55, 5.4, 3.0, [
    [seg("Raise the skill floor of agentic coding — so being a great agent-wrangler stops being the moat, and product taste is the only thing that still matters.", TEXT, 15, font=BODY)],
    [seg("", TEXT, 8)],
    [seg("For my own work, I've mostly already done that with Smithers.", DIM, 13.5, font=BODY)],
    [seg("", TEXT, 6)],
    [seg("The company is that personal result, productized.", GREEN, 14.5, bold=True, font=BODY)],
], space_after=4, line_spacing=1.16)
notes(s, "Let me talk about why I'm the person to build this. I'm a highly respected agentic coder. "
         "I've shipped serious production code using models far weaker than today's — including the "
         "fastest Ethereum virtual machine ever written, faster than the previous leader, revm, and I "
         "built it with Claude three-point-five. I also spent time at Google, and a lot of how "
         "Smithers organizes large numbers of agents working in parallel is taken directly from how "
         "Google organizes its monorepo. Here's the thesis: I believe I can raise the skill floor of "
         "agentic coders, so that being an elite agent-wrangler stops being the differentiator, and "
         "the only thing that still matters is being a top-tier product developer. For my own work, "
         "I've mostly already done that with Smithers. The company is that personal result, "
         "productized.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — GO-TO-MARKET
# ════════════════════════════════════════════════════════════════════════════
s = base("GO-TO-MARKET", "Go-to-market — reputation-led",
         "Land enterprises through relationships and credibility, not cold sales.")
steps = [("1", "Best-in-class product", "we stand behind it · OSS credibility de-risks it"),
         ("2", "Founder relationships", "trust + standing in the industry opens the door"),
         ("3", "Hands-on onboarding", "walk the org onto real workflows, not a trial"),
         ("4", "Land + expand", "workflows go load-bearing → enterprise contract")]
sx = 0.6; sw = 2.95; sg = 0.12
for i, (n, t1, t2) in enumerate(steps):
    x = sx + i * (sw + sg)
    card(s, x, 2.85, sw, 2.05, fill=PANEL, line=BORDER, radius=0.06)
    txt(s, x + 0.25, 3.05, 1.0, 0.6, [[seg(n, CYAN, 26, bold=True, font=MONO)]])
    txt(s, x + 0.25, 3.7, sw - 0.5, 0.5, [[seg(t1, WHITE, 13.5, bold=True, font=HEAD)]], line_spacing=1.0)
    txt(s, x + 0.25, 4.18, sw - 0.5, 0.6, [[seg(t2, DIM, 11, font=BODY)]], line_spacing=1.1)
    if i < 3:
        arrow_glyph(s, x + sw, 2.85, sg, color=MUTED, size=16)
card(s, 0.6, 5.15, 12.13, 1.6, fill=rgb("11202B"), line=YELLOW, radius=0.05)
txt(s, 0.95, 5.36, 11.5, 0.4, [[seg("In progress: early conversations with Opendoor", YELLOW, 16, bold=True, font=HEAD)]])
txt(s, 0.95, 5.82, 11.5, 0.4, [[seg("re: a consulting engagement to onboard their org onto Smithers — a live example of the motion, not a closed deal.", DIM, 13, font=BODY)]], line_spacing=1.1)
txt(s, 0.95, 6.3, 11.5, 0.35, [[seg("Reputation-based onboarding: lower CAC, higher trust, stickier landings.", TEAL, 12.5, italic=True, font=BODY)]])
notes(s, "On go-to-market. The product lands bottom-up through open source. But for enterprise, our "
         "wedge is reputation. We have a product we genuinely stand by as best in class, and we "
         "intend to leverage my relationships and standing in the industry to do reputation-based "
         "onboarding. Concretely: we're in early conversations with Opendoor about bringing me on as "
         "a consultant to help onboard their organization onto Smithers. That's the model — a trusted "
         "operator walks an org onto the tool, the workflows become load-bearing, and the enterprise "
         "relationship grows from there. It's far cheaper and far higher-trust than a traditional "
         "sales motion.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — WHY NOW, WHY US
# ════════════════════════════════════════════════════════════════════════════
s = base("THE ASK", "Why now, why us",
         "The timing is forced. The wedge is open. The data compounds.")
points = [
    (GREEN, "Forced timing", "agents got good; the reliability layer to run them like a business doesn't exist yet"),
    (CYAN, "Real traction", "21,001 downloads/mo and accelerating · 268 stars · 15+ external contributors"),
    (YELLOW, "Live ecosystem", "third-party UIs + Kubernetes automation feeding the roadmap"),
    (MAGENTA, "Compounding moat", "OSS distribution + switching cost per workflow landed"),
    (TEAL, "Reputation GTM", "relationship-led enterprise landings (e.g. Opendoor, in progress)"),
]
py = 2.8
for col, k, v in points:
    card(s, 0.6, py, 12.13, 0.78, fill=PANEL, line=None, radius=0.06)
    rect(s, 0.6, py, 0.1, 0.78, fill=col)
    txt(s, 0.95, py, 3.4, 0.78, [[seg("▸ ", col, 14, bold=True), seg(k, WHITE, 15.5, bold=True, font=HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 4.35, py, 8.2, 0.78, [[seg(v, DIM, 13, font=BODY)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    py += 0.88
notes(s, "Why now and why us. Why now: agents just crossed the line from demo to dependable, and the "
         "reliability layer to run them like a business does not exist yet. That window is open right "
         "now. Why us: we already have the adoption, the contributor base, a real third-party "
         "ecosystem feeding our roadmap, and a reputation-led path into the enterprise that most "
         "infrastructure startups would kill for. We're early, the numbers are accelerating, and the "
         "moat — open-source distribution plus switching cost — compounds with every workflow that "
         "moves onto us.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — LET'S TALK (close)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); _idx += 1
rect(s, 0, 0, SW, SH, fill=BG)
place_image(s, "cli-hero.png", 0, 0, SW, SH, frame=False)
# darken overlay
ov = rect(s, 0, 0, SW, SH, fill=rgb("0A0E16"))
ov.fill.fore_color.rgb = rgb("0A0E16")
sp = ov.fill._xPr.find(qn('a:solidFill'))
clr = sp.find(qn('a:srgbClr'))
a = clr.makeelement(qn('a:alpha'), {'val': '62000'}); clr.append(a)
s.shapes.add_picture(os.path.join(ASSETS, "logo.png"), In(5.55), In(1.55), In(0.95), In(0.95))
txt(s, 6.6, 1.7, 4, 0.7, [[seg("smithers", WHITE, 30, bold=True, font=MONO)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 1.0, 2.95, 11.33, 0.6, [[seg("The durable control plane for humans + agents doing real work.", WHITE, 23, bold=True, font=HEAD)]], align=PP_ALIGN.CENTER)
chips = [("36,913+", "downloads", GREEN), ("268", "stars", YELLOW), ("15+", "contributors", CYAN), ("21,001", "/mo, accelerating", TEAL)]
cw, gap = 2.2, 0.2
total_w = len(chips) * cw + (len(chips) - 1) * gap
x0 = (SW - total_w) / 2
for i, (num, lab, col) in enumerate(chips):
    x = x0 + i * (cw + gap)
    card(s, x, 3.95, cw, 1.0, fill=rgb("12182E"), line=BORDER, radius=0.08)
    txt(s, x, 4.08, cw, 0.45, [[seg(num, col, 22, bold=True, font=MONO)]], align=PP_ALIGN.CENTER)
    txt(s, x, 4.62, cw, 0.3, [[seg(lab, DIM, 11, font=BODY)]], align=PP_ALIGN.CENTER)
txt(s, 1.0, 5.45, 11.33, 0.4, [[seg("smithers.sh", WHITE, 16, bold=True, font=MONO),
                                seg("   ·   github.com/smithersai/smithers", DIM, 16, font=MONO)]], align=PP_ALIGN.CENTER)
txt(s, 1.0, 6.25, 11.33, 0.5, [[seg("Let's talk.", WHITE, 26, bold=True, font=HEAD)]], align=PP_ALIGN.CENTER)
notes(s, "That's Smithers. The durable control plane for humans and agents doing real work, with real "
         "adoption, a turning flywheel, and a credible path into the enterprise. I'd love to tell you "
         "more, and I want to hear where you'd push on it. Thank you.")

prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
